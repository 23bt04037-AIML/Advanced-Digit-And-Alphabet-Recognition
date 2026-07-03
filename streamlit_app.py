"""
DigitAI – Advanced Handwritten Digit Recognition Using CNN
Full Streamlit App  |  Landing → Login → Role-Based Dashboard
Features: Batch Prediction · Forgot Password (OTP) · Logout Confirm · MySQL ready
"""
import os, sys, json, uuid, io, re, random, smtplib
import datetime as _dt
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from pathlib import Path
import numpy as np
import pymysql

ROOT = Path(__file__).parent
sys.path.insert(0, str(ROOT))
os.chdir(ROOT)

# Load SMTP/email settings from .env file for OTP email sending
try:
    from dotenv import load_dotenv
    load_dotenv(ROOT / ".env")
except Exception:
    # App can still run; OTP email will show a clear SMTP warning if .env is missing.
    pass


import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from PIL import Image, ImageOps
import streamlit as st
from streamlit_drawable_canvas import st_canvas
import uuid


# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Advanced Digit Recognition – CNN",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ══════════════════════════════════════════════════════════════════════════════
# MODERN THEME + GLOBAL CSS
# ══════════════════════════════════════════════════════════════════════════════
if "dark_mode" not in st.session_state:
    st.session_state["dark_mode"] = True


def _theme_palette():
    """Return color tokens for light/dark UI."""
    dark = bool(st.session_state.get("dark_mode", True))
    if dark:
        return {
            "name": "Dark",
            "bg": "#070A13",
            "bg2": "#0B1020",
            "surface": "rgba(15, 23, 42, 0.88)",
            "surface2": "rgba(30, 41, 59, 0.70)",
            "card": "rgba(15, 23, 42, 0.78)",
            "text": "#E5E7EB",
            "muted": "#94A3B8",
            "soft": "#CBD5E1",
            "primary": "#60A5FA",
            "primary2": "#8B5CF6",
            "accent": "#22C55E",
            "warning": "#F59E0B",
            "danger": "#EF4444",
            "border": "rgba(148, 163, 184, 0.22)",
            "input": "rgba(2, 6, 23, 0.55)",
            "shadow": "0 20px 70px rgba(0,0,0,.40)",
            "hero": "radial-gradient(circle at 15% 20%, rgba(96,165,250,.32), transparent 30%), radial-gradient(circle at 85% 10%, rgba(139,92,246,.32), transparent 26%), linear-gradient(135deg, #0B1020 0%, #111827 48%, #1E1B4B 100%)",
            "hero_text": "#F8FAFC",
            "canvas_border": "#60A5FA",
        }
    return {
        "name": "Light",
        "bg": "#F8FAFC",
        "bg2": "#EEF2FF",
        "surface": "rgba(255, 255, 255, 0.92)",
        "surface2": "rgba(248, 250, 252, 0.90)",
        "card": "rgba(255, 255, 255, 0.86)",
        "text": "#0F172A",
        "muted": "#475569",
        "soft": "#334155",
        "primary": "#2563EB",
        "primary2": "#7C3AED",
        "accent": "#16A34A",
        "warning": "#D97706",
        "danger": "#DC2626",
        "border": "rgba(15, 23, 42, 0.12)",
        "input": "rgba(255, 255, 255, 0.92)",
        "shadow": "0 20px 70px rgba(15,23,42,.12)",
        "hero": "radial-gradient(circle at 15% 20%, rgba(37,99,235,.18), transparent 30%), radial-gradient(circle at 85% 10%, rgba(124,58,237,.16), transparent 26%), linear-gradient(135deg, #FFFFFF 0%, #EFF6FF 48%, #EEF2FF 100%)",
        "hero_text": "#0F172A",
        "canvas_border": "#2563EB",
    }


def _apply_theme_css():
    t = _theme_palette()
    st.markdown(f"""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&display=swap');

:root {{
  --bg: {t['bg']};
  --bg2: {t['bg2']};
  --surface: {t['surface']};
  --surface2: {t['surface2']};
  --card: {t['card']};
  --text: {t['text']};
  --muted: {t['muted']};
  --soft: {t['soft']};
  --primary: {t['primary']};
  --primary2: {t['primary2']};
  --accent: {t['accent']};
  --warning: {t['warning']};
  --danger: {t['danger']};
  --border: {t['border']};
  --input: {t['input']};
  --shadow: {t['shadow']};
  --hero: {t['hero']};
  --hero-text: {t['hero_text']};
  --canvas-border: {t['canvas_border']};
}}

html, body, [class*="css"] {{ font-family: 'Inter', sans-serif; }}
#MainMenu {{ visibility:hidden; }}
footer {{ visibility:hidden; }}
header {{ background: transparent !important; }}

.stApp {{
  background:
    radial-gradient(circle at top left, rgba(96,165,250,.14), transparent 28%),
    radial-gradient(circle at bottom right, rgba(139,92,246,.14), transparent 30%),
    linear-gradient(180deg, var(--bg), var(--bg2));
  color: var(--text);
}}

.block-container {{
  padding-top: 1.1rem !important;
  padding-bottom: 3rem !important;
  max-width: 1500px;
}}

h1, h2, h3, h4, h5, h6, p, label, span, div {{ color: inherit; }}
h1, h2, h3 {{ color: var(--text) !important; letter-spacing: -0.03em; }}
p, label, .stMarkdown, .stCaption {{ color: var(--muted); }}

/* Top bars */
.app-topbar {{
  background: var(--surface);
  border: 1px solid var(--border);
  box-shadow: var(--shadow);
  backdrop-filter: blur(18px);
  border-radius: 22px;
  padding: .85rem 1.15rem;
  margin: .2rem 0 1rem 0;
  display: flex;
  align-items: center;
  justify-content: space-between;
  gap: 1rem;
}}
.brand-wrap {{display:flex; align-items:center; gap:.9rem;}}
.brand-logo {{
  width:42px; height:42px; border-radius:14px;
  display:grid; place-items:center;
  background: linear-gradient(135deg, var(--primary), var(--primary2));
  box-shadow: 0 12px 28px rgba(37,99,235,.28);
  font-size:1.35rem;
}}
.brand-title {{font-size:1.12rem; font-weight:900; color:var(--text); line-height:1;}}
.brand-subtitle {{font-size:.80rem; color:var(--muted); margin-top:.18rem;}}
.topbar-user {{color:var(--soft); font-size:.9rem; font-weight:600;}}

/* Hero */
.hero-shell {{
  background: var(--hero);
  color: var(--hero-text);
  border: 1px solid var(--border);
  border-radius: 28px;
  padding: 3.2rem 2.2rem;
  box-shadow: var(--shadow);
  overflow: hidden;
  position: relative;
}}
.hero-shell:after {{
  content:""; position:absolute; inset:auto -10% -35% auto;
  width:360px; height:360px; border-radius:50%;
  background: linear-gradient(135deg, rgba(96,165,250,.20), rgba(139,92,246,.18));
  filter: blur(8px);
}}
.hero-kicker {{
  display:inline-flex; align-items:center; gap:.45rem;
  padding:.40rem .75rem; border:1px solid var(--border);
  border-radius:999px; background:rgba(255,255,255,.08);
  color:var(--hero-text); font-weight:700; font-size:.8rem;
}}
.hero-title {{
  color:var(--hero-text) !important;
  font-size: clamp(2.1rem, 4vw, 4.4rem);
  line-height: 1.02;
  margin: 1rem 0 .8rem 0;
  font-weight: 900;
  max-width: 950px;
}}
.hero-text {{
  color:var(--muted);
  max-width: 850px;
  font-size:1.02rem;
  line-height:1.8;
}}

.feature-grid {{display:grid; grid-template-columns:repeat(3,minmax(0,1fr)); gap:1rem; margin-top:1.2rem;}}
@media (max-width: 900px) {{ .feature-grid {{grid-template-columns:1fr;}} }}
.feature-card, .modern-card, .auth-card {{
  background: var(--card);
  border: 1px solid var(--border);
  box-shadow: var(--shadow);
  backdrop-filter: blur(18px);
  border-radius: 24px;
  padding: 1.35rem;
}}
.feature-card {{transition: transform .2s ease, border-color .2s ease;}}
.feature-card:hover {{transform: translateY(-4px); border-color: rgba(96,165,250,.55);}}
.feature-card h3, .feature-card h4 {{color:var(--text) !important; margin:.35rem 0;}}
.feature-card p {{color:var(--muted); margin:0; line-height:1.6;}}
.feature-icon {{
  width:46px; height:46px; display:grid; place-items:center; border-radius:16px;
  background: linear-gradient(135deg, var(--primary), var(--primary2));
  color:white; font-size:1.25rem; box-shadow: 0 12px 28px rgba(37,99,235,.25);
}}

.page-hero {{
  background: var(--hero);
  border:1px solid var(--border);
  border-radius: 26px;
  padding: 1.5rem 1.6rem;
  margin: .4rem 0 1.2rem 0;
  box-shadow: var(--shadow);
}}
.page-hero h2 {{margin:0 0 .35rem 0; color: var(--hero-text) !important;}}
.page-hero p {{margin:0; color:var(--muted); line-height:1.65;}}
.mode-chip {{
  display:inline-flex; align-items:center; gap:.4rem;
  padding:.35rem .7rem; border-radius:999px;
  background:rgba(96,165,250,.12);
  border:1px solid var(--border); color:var(--primary); font-weight:800; font-size:.78rem;
}}

/* Streamlit widgets */
.stButton > button, .stDownloadButton > button, button[kind="primary"] {{
  border-radius: 14px !important;
  border: 1px solid var(--border) !important;
  font-weight: 800 !important;
  min-height: 2.75rem;
  box-shadow: 0 10px 24px rgba(37,99,235,.16);
}}
.stButton > button[kind="primary"], .stFormSubmitButton > button[kind="primary"] {{
  background: linear-gradient(135deg, var(--primary), var(--primary2)) !important;
  color: white !important;
  border: none !important;
}}
.stButton > button:hover, .stDownloadButton > button:hover {{
  transform: translateY(-1px);
  border-color: var(--primary) !important;
}}

div[data-baseweb="select"] > div,
.stTextInput input,
.stNumberInput input,
.stTextArea textarea,
.stFileUploader section,
.stSlider,
.stCheckbox,
.stRadio,
.stMultiSelect,
.stDateInput input {{
  border-radius: 16px !important;
}}
.stTextInput input, .stNumberInput input, .stTextArea textarea {{
  background: var(--input) !important;
  color: var(--text) !important;
  border: 1px solid var(--border) !important;
}}
.stFileUploader section {{
  background: var(--card) !important;
  border: 1.5px dashed var(--primary) !important;
}}

.stTabs [data-baseweb="tab-list"] {{
  gap: .45rem;
  background: var(--surface);
  border:1px solid var(--border);
  border-radius: 20px;
  padding: .45rem;
  box-shadow: var(--shadow);
}}
.stTabs [data-baseweb="tab"] {{
  border-radius: 14px;
  padding: .55rem 1rem;
  color: var(--muted);
  font-weight: 800;
}}
.stTabs [aria-selected="true"] {{
  background: linear-gradient(135deg, var(--primary), var(--primary2));
  color: white !important;
}}

[data-testid="stMetric"] {{
  background: var(--card);
  border: 1px solid var(--border);
  border-radius: 18px;
  padding: 1rem;
  box-shadow: 0 10px 30px rgba(15,23,42,.10);
}}
[data-testid="stMetricValue"] {{color: var(--text) !important; font-weight:900;}}
[data-testid="stMetricLabel"] {{color: var(--muted) !important;}}

/* Prediction result */
.pred-box {{
  background: var(--card);
  border: 1px solid var(--border);
  box-shadow: var(--shadow);
  border-radius: 24px;
  padding: 1.6rem;
  text-align:center;
  margin-bottom:1rem;
}}
.pred-digit {{font-size:5.4rem;font-weight:900;color:var(--primary);line-height:1;}}
.pred-label {{color:var(--muted);font-size:.9rem;margin-top:.3rem;font-weight:700;}}
.conf-pill {{display:inline-block;padding:.35rem 1.1rem;border-radius:999px;font-weight:800;font-size:.95rem;margin-top:.7rem;}}
.high {{background:rgba(34,197,94,.14);color:#22C55E;border:1px solid rgba(34,197,94,.38);}}
.mid  {{background:rgba(245,158,11,.14);color:#F59E0B;border:1px solid rgba(245,158,11,.38);}}
.low  {{background:rgba(239,68,68,.14);color:#EF4444;border:1px solid rgba(239,68,68,.38);}}

.section-title {{
  font-size:.98rem;font-weight:900;color:var(--text);
  border-left:4px solid var(--primary);
  padding-left:.75rem;margin:1.25rem 0 .65rem;
}}
.role-admin, .role-client {{padding:3px 11px;border-radius:99px;font-size:.72rem;font-weight:900;color:white;}}
.role-admin {{background:linear-gradient(135deg,#7c3aed,#db2777);}}
.role-client {{background:linear-gradient(135deg,#2563eb,#06b6d4);}}

[data-testid="stCanvas"] {{
  border:2px solid var(--canvas-border) !important;
  border-radius:18px !important;
  box-shadow: 0 12px 30px rgba(37,99,235,.16);
  overflow:hidden;
}}

[data-testid="stDataFrame"] {{border-radius:18px; overflow:hidden; border:1px solid var(--border);}}
.streamlit-expanderHeader {{background:var(--card); border-radius:14px;}}
hr {{border-color: var(--border) !important;}}
code {{border-radius: 14px !important;}}


/* Compact + reliable light/dark fixes */
.stApp, .stApp * {{ transition: background-color .18s ease, color .18s ease, border-color .18s ease; }}

/* Fix old hardcoded dark auth blocks so text/background follows theme */
div[style*="background:#1a1a2e"],
div[style*="background: #1a1a2e"],
div[style*="background:#0f0f1a"],
div[style*="background:linear-gradient(90deg,#0f0f1a,#1a1a2e)"] {{
  background: var(--card) !important;
  border-color: var(--border) !important;
  box-shadow: var(--shadow) !important;
  color: var(--text) !important;
}}
div[style*="background:#1a1a2e"] h1,
div[style*="background:#1a1a2e"] h2,
div[style*="background:#1a1a2e"] h3,
div[style*="background:#1a1a2e"] p,
div[style*="background:#1a1a2e"] span {{
  color: var(--text) !important;
}}

/* Streamlit forms should not look detached */
[data-testid="stForm"] {{
  background: var(--surface) !important;
  border: 1px solid var(--border) !important;
  border-radius: 22px !important;
  padding: 1.1rem !important;
  box-shadow: 0 14px 40px rgba(15,23,42,.10) !important;
}}

/* Smaller page headers */
.page-hero {{
  padding: 1rem 1.25rem !important;
  border-radius: 20px !important;
  margin: .35rem 0 1rem 0 !important;
}}
.page-hero h2 {{ font-size: 1.85rem !important; margin: .35rem 0 .25rem 0 !important; }}
.page-hero p {{ font-size: .95rem !important; line-height: 1.5 !important; }}
.mode-chip {{ padding: .25rem .55rem !important; font-size: .72rem !important; }}

/* Cleaner top spacing */
.block-container {{ padding-top: .75rem !important; }}

/* Keep radios/selects readable in both themes */
[data-testid="stRadio"] label, [data-testid="stRadio"] p,
[data-testid="stSelectbox"] label, [data-testid="stFileUploader"] label {{
  color: var(--text) !important;
}}

/* Reduce visual overload: hide long guide cards */
.clear-guide, .clear-steps, .choice-help {{
  display: none !important;
}}
.clear-mini-title {{
  margin: .35rem 0 .15rem 0 !important;
  font-size: 1rem !important;
}}
.clear-mini-text {{
  margin-bottom: .55rem !important;
  font-size: .88rem !important;
}}
</style>
""", unsafe_allow_html=True)


_apply_theme_css()


# Extra clarity UI layer
st.markdown("""
<style>
.clear-guide {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 22px;
  padding: 1.05rem 1.15rem;
  margin: .7rem 0 1rem 0;
  box-shadow: 0 12px 35px rgba(15,23,42,.10);
}
.clear-guide h3 {
  margin: 0 0 .35rem 0 !important;
  color: var(--text) !important;
  font-size: 1.05rem !important;
}
.clear-guide p { margin: 0; color: var(--muted); line-height: 1.65; }
.clear-steps {
  display: grid;
  grid-template-columns: repeat(4, minmax(0, 1fr));
  gap: .8rem;
  margin: .9rem 0 1.1rem 0;
}
@media (max-width: 1000px) { .clear-steps { grid-template-columns: repeat(2, minmax(0, 1fr)); } }
@media (max-width: 650px) { .clear-steps { grid-template-columns: 1fr; } }
.clear-step {
  background: var(--card);
  border: 1px solid var(--border);
  border-radius: 18px;
  padding: .95rem;
  min-height: 108px;
}
.clear-step .num {
  width: 28px;
  height: 28px;
  display: inline-grid;
  place-items: center;
  border-radius: 999px;
  background: linear-gradient(135deg, var(--primary), var(--primary2));
  color: white;
  font-weight: 900;
  margin-bottom: .55rem;
}
.clear-step b { color: var(--text); display: block; margin-bottom: .2rem; }
.clear-step span { color: var(--muted); font-size: .86rem; line-height: 1.45; }
.choice-help {
  display: grid;
  grid-template-columns: repeat(4, minmax(0, 1fr));
  gap: .8rem;
  margin: .75rem 0 1rem 0;
}
@media (max-width: 1150px) { .choice-help { grid-template-columns: repeat(2, minmax(0, 1fr)); } }
@media (max-width: 700px) { .choice-help { grid-template-columns: 1fr; } }
.choice-card {
  background: var(--card);
  border: 1px solid var(--border);
  border-radius: 18px;
  padding: .95rem;
}
.choice-card.active {
  border-color: var(--primary);
  box-shadow: 0 0 0 3px rgba(96,165,250,.16), 0 12px 32px rgba(15,23,42,.12);
}
.choice-card h4 { margin: 0 0 .25rem 0 !important; color: var(--text) !important; font-size: .98rem; }
.choice-card p { margin: 0; color: var(--muted); font-size: .86rem; line-height: 1.55; }
.clear-mini-title {
  font-size: 1.1rem;
  font-weight: 900;
  color: var(--text);
  margin: .8rem 0 .15rem 0;
}
.clear-mini-text {
  color: var(--muted);
  margin-bottom: .8rem;
  line-height: 1.6;
}
.result-placeholder {
  background: var(--card);
  border: 1px dashed var(--primary);
  border-radius: 20px;
  padding: 1.2rem;
  color: var(--muted);
  line-height: 1.7;
}
[data-testid="stRadio"] label p { font-weight: 800 !important; }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# CONSTANTS / LAZY LOADERS
# ══════════════════════════════════════════════════════════════════════════════
MODELS_DIR = ROOT / "models"
# Models that are NOT digit classifiers must be excluded from the digit model selector.
# This prevents OCR/CTC models such as handwriting_ocr_prediction_model.keras
# from appearing in digit model dropdowns.
_DIGIT_MODEL_EXCLUDE = {
    "ocr_prediction_model",
    "ocr_training_model",
    "best_ocr_training_model",
    "best_ocr_prediction_model",
    "handwriting_ocr_prediction_model",
    "handwriting_ocr_training_model",
    "best_handwriting_ocr_model",
    "handwriting_ocr",
    "best_digit_model",
}

def _base_model_name(stem: str) -> str:
    return (
        stem
        .split("_adam")[0]
        .split("_nadam")[0]
        .split("_rmsprop")[0]
        .split("_sgd")[0]
    )

def _is_digit_model_name(model_name: str) -> bool:
    name = model_name.lower()

    blocked_keywords = [
        "ocr",
        "ctc",
        "charset",
        "handwriting",
        "text",
        "word",
        "alphabet",
        "script",
    ]

    if model_name in _DIGIT_MODEL_EXCLUDE:
        return False

    if any(k in name for k in blocked_keywords):
        return False

    allowed_keywords = [
        "cnn",
        "mobilenet",
        "resnet",
        "digit",
    ]

    return any(k in name for k in allowed_keywords)

AVAILABLE_MODELS = sorted(
    {
        _base_model_name(p.stem)
        for p in MODELS_DIR.glob("*.keras")
        if not p.stem.endswith("_best")
        and _is_digit_model_name(_base_model_name(p.stem))
    }
) or ["cnn_medium"]

MODEL_LABELS = {
    "cnn_small":   "CNN Small  (lightweight)",
    "cnn_medium":  "CNN Medium (balanced) ✦",
    "cnn_deep":    "CNN Deep   (highest accuracy)",
    "mobilenetv2": "MobileNetV2 (transfer learning)",
    "resnet50":    "ResNet50    (transfer learning)",
}

_IST = _dt.timezone(_dt.timedelta(hours=5, minutes=30))

def _to_ist(ts):
    if ts is None: return ""
    if isinstance(ts, str):
        try: ts = _dt.datetime.fromisoformat(ts)
        except: return ts[:16]
    # The database timestamps appear to already be in local time on this system,
    # so we don't need to add another +05:30.
    return ts.strftime("%Y-%m-%d %H:%M")


@st.cache_resource(show_spinner="Loading prediction engine…")
def load_prediction_service():
    from backend.services.prediction_service import predict, generate_gradcam, generate_lime, generate_fgsm_attack, registry
    return predict, generate_gradcam, generate_lime, generate_fgsm_attack, registry


@st.cache_resource(show_spinner="Connecting to database…")
def load_db():
    from backend.database.db import SessionLocal, create_tables, check_connection
    check_connection(); create_tables(); _ensure_admin_exists()
    return SessionLocal


def _ensure_admin_exists():
    try:
        from backend.database.db import SessionLocal
        from backend.database.models import User
        db = SessionLocal()
        if db.query(User).filter(User.role == "admin").count() == 0:
            admin = User(username="admin", password_hash=_hash_pw("Admin@123"),
                         role="admin", full_name="Administrator",
                         email="admin@digitai.local")
            db.add(admin); db.commit()
        db.close()
    except: pass


# ══════════════════════════════════════════════════════════════════════════════
# AUTH HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def _hash_pw(pw: str) -> str:
    import hashlib
    return hashlib.sha256(pw.encode()).hexdigest()

def _validate_password(pw: str):
    if len(pw) < 8:                       return "Min 8 characters required."
    if not re.search(r"[A-Z]", pw):       return "At least 1 uppercase letter required."
    if not re.search(r"[^A-Za-z0-9]", pw):return "At least 1 special character (!@#$% …) required."
    return None

def _login(username: str, password: str):
    try:
        SL = load_db()
        from backend.database.models import User
        db = SL()
        u = db.query(User).filter(User.username == username, User.is_active == True).first()
        if not u:             db.close(); return None, "Username not found."
        if u.password_hash != _hash_pw(password):
            db.close(); return None, "Incorrect password."
        u.last_login = _dt.datetime.utcnow(); db.commit()
        result = {"id": u.id, "username": u.username,
                  "role": u.role, "full_name": u.full_name or u.username,
                  "email": u.email or ""}
        db.close(); return result, None
    except Exception as e:
        return None, f"Login error: {e}"

def _register(username, password, full_name, email):
    """Create new client account and return user data."""
    err = _validate_password(password)
    if err:
        return None, err

    try:
        SL = load_db()
        from backend.database.models import User

        db = SL()

        if db.query(User).filter(User.username == username).first():
            db.close()
            return None, "Username already taken."

        if db.query(User).filter(User.email == email).first():
            db.close()
            return None, "Email already registered."

        new_user = User(
            username=username,
            password_hash=_hash_pw(password),
            role="client",
            full_name=full_name,
            email=email,
            is_active=True,
            last_login=_dt.datetime.utcnow()
        )

        db.add(new_user)
        db.commit()
        db.refresh(new_user)

        user_data = {
            "id": new_user.id,
            "username": new_user.username,
            "role": new_user.role,
            "full_name": new_user.full_name or new_user.username,
            "email": new_user.email or "",
        }

        db.close()
        return user_data, None

    except Exception as e:
        return None, f"Registration error: {e}"


# ══════════════════════════════════════════════════════════════════════════════
# OTP / FORGOT PASSWORD HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _get_config_value(name: str, default: str = "") -> str:
    """
    Read config from .env/environment first, then Streamlit secrets.
    This helps OTP email work in both local system and Streamlit deployment.
    """
    value = os.getenv(name)
    if value is not None and str(value).strip() != "":
        return str(value).strip()
    try:
        value = st.secrets.get(name, default)
        return str(value).strip() if value is not None else default
    except Exception:
        return default


def _send_email(to_addr: str, subject: str, html_body: str) -> str | None:
    """
    Send OTP/reset emails using SMTP.

    For Gmail:
    SMTP_HOST=smtp.gmail.com
    SMTP_PORT=587
    SMTP_USER=yourgmail@gmail.com
    SMTP_PASSWORD=your 16-character Gmail App Password
    SMTP_FROM=yourgmail@gmail.com
    """
    host = _get_config_value("SMTP_HOST", "")
    port_raw = _get_config_value("SMTP_PORT", "587")
    user = _get_config_value("SMTP_USER", "")
    pwd = _get_config_value("SMTP_PASSWORD", "")
    frm = _get_config_value("SMTP_FROM", user)

    try:
        port = int(port_raw)
    except Exception:
        port = 587

    # Gmail app passwords are sometimes copied with spaces; remove them.
    pwd = str(pwd).replace(" ", "")

    if not host or not user or not pwd:
        return (
            "SMTP not configured. Create a .env file with SMTP_HOST, SMTP_PORT, "
            "SMTP_USER, SMTP_PASSWORD, and SMTP_FROM."
        )

    try:
        msg = MIMEMultipart("alternative")
        msg["Subject"] = subject
        msg["From"] = frm or user
        msg["To"] = to_addr
        msg.attach(MIMEText(html_body, "html"))

        use_ssl = port == 465 or _get_config_value("SMTP_SSL", "false").lower() in {"1", "true", "yes"}

        if use_ssl:
            with smtplib.SMTP_SSL(host, port, timeout=20) as s:
                s.login(user, pwd)
                s.sendmail(frm or user, [to_addr], msg.as_string())
        else:
            with smtplib.SMTP(host, port, timeout=20) as s:
                s.ehlo()
                s.starttls()
                s.ehlo()
                s.login(user, pwd)
                s.sendmail(frm or user, [to_addr], msg.as_string())

        return None

    except smtplib.SMTPAuthenticationError:
        return (
            "SMTP authentication failed. For Gmail, use a Gmail App Password, "
            "not your normal Gmail password."
        )
    except smtplib.SMTPConnectError:
        return "Could not connect to SMTP server. Check SMTP_HOST, SMTP_PORT, and internet connection."
    except Exception as e:
        return f"SMTP error: {e}"


def _generate_otp() -> str:
    return str(random.randint(100000, 999999))

def _store_otp(user_id: int, otp: str) -> bool:
    try:
        SL = load_db()
        from backend.database.models import OTPRecord
        db = SL()
        # invalidate old OTPs
        db.query(OTPRecord).filter(OTPRecord.user_id == user_id, OTPRecord.used == False).update({"used": True})
        expires = _dt.datetime.utcnow() + _dt.timedelta(minutes=10)
        db.add(OTPRecord(user_id=user_id, otp_code=otp, expires_at=expires))
        db.commit(); db.close(); return True
    except: return False

def _verify_otp(user_id: int, otp: str) -> str | None:
    """Return None if valid, error string otherwise."""
    try:
        SL = load_db()
        from backend.database.models import OTPRecord
        db = SL()
        record = (db.query(OTPRecord)
                  .filter(OTPRecord.user_id == user_id,
                          OTPRecord.otp_code == otp,
                          OTPRecord.used == False)
                  .order_by(OTPRecord.created_at.desc()).first())
        if not record:
            db.close(); return "Invalid OTP."
        if _dt.datetime.utcnow() > record.expires_at:
            db.close(); return "OTP expired. Please request a new one."
        record.used = True; db.commit(); db.close(); return None
    except Exception as e:
        return f"OTP verify error: {e}"

def _reset_password(user_id: int, new_pw: str) -> str | None:
    try:
        SL = load_db()
        from backend.database.models import User
        db = SL()
        u = db.query(User).filter(User.id == user_id).first()
        if not u: db.close(); return "User not found."
        u.password_hash = _hash_pw(new_pw); db.commit()
        email = u.email; db.close(); return email
    except Exception as e:
        return None


# ══════════════════════════════════════════════════════════════════════════════
# LANDING PAGE  (before login – project info only)
# ══════════════════════════════════════════════════════════════════════════════
def page_landing():
    _auth_topbar()

    st.markdown("""
    <div class="hero-shell">
      <div class="hero-kicker">✨ AI Powered Recognition System</div>
      <h1 class="hero-title">Advanced  Digit, Multi-Digit & OCR Recognition</h1>
      <p class="hero-text">
        DigitAI combines CNN-based digit recognition, universal multi-digit detection,
        OCR for alphabets/words, batch processing, PDF/Word/PPT text extraction,
        analytics dashboards, and user history in one modern interface.
      </p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div class="feature-grid">
      <div class="feature-card">
        <div class="feature-icon">🔢</div>
        <h3>Smart Digit AI</h3>
        <p>Recognize single digits, multi-digit numbers, canvas drawings, and uploaded images.</p>
      </div>
      <div class="feature-card">
        <div class="feature-icon">📄</div>
        <h3>OCR + Documents</h3>
        <p>Extract text from images, PDF pages, Word documents, and PowerPoint slides.</p>
      </div>
      <div class="feature-card">
        <div class="feature-icon">📊</div>
        <h3>Analytics Dashboard</h3>
        <p>Track prediction history, confidence, model performance, and export reports.</p>
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 1.1, 1])
    with c2:
        st.markdown("""
        <div class="modern-card" style="text-align:center;">
          <div class="mode-chip">🚀 Start Recognition</div>
          <h3 style="margin:.75rem 0 .25rem 0;">Continue to DigitAI</h3>
          <p style="margin:0;">Login or create a new account to use the full system.</p>
        </div>
        """, unsafe_allow_html=True)
        st.markdown("<br>", unsafe_allow_html=True)
        b1, b2 = st.columns(2)
        with b1:
            if st.button("🔑 Login", type="primary", use_container_width=True):
                st.session_state["nav"] = "login"
                st.rerun()
        with b2:
            if st.button("📝 Register", use_container_width=True):
                st.session_state["nav"] = "register"
                st.rerun()


def page_login():
    _auth_topbar()
    _, mid, _ = st.columns([1, 1.4, 1])

    with mid:
        st.markdown("""
        <div class="auth-card" style="text-align:center;margin-top:1rem;">
          <div style="font-size:2rem">🔑</div>
          <h2 style="margin:.35rem 0 .15rem 0;font-size:1.45rem">Sign In</h2>
          <p style="margin:0">Enter your username and password.</p>
        </div>
        """, unsafe_allow_html=True)

        with st.form("login_form", clear_on_submit=False):
            uname = st.text_input("Username", placeholder="Enter username")
            pwd = st.text_input("Password", type="password", placeholder="Enter password")
            ok = st.form_submit_button("Login →", type="primary", use_container_width=True)

            if ok:
                if not uname or not pwd:
                    st.error("Please enter both fields.")
                else:
                    user, err = _login(uname, pwd)

                    if err:
                        st.error(f"❌ {err}")
                    else:
                        st.session_state["auth_user"] = user
                        st.session_state["session_id"] = uuid.uuid4().hex
                        st.session_state.pop("nav", None)
                        st.rerun()

        c1, c2, c3 = st.columns(3)

        with c1:
            if st.button("← Back", use_container_width=True):
                st.session_state["nav"] = "landing"
                st.rerun()

        with c2:
            if st.button("New User / Register", use_container_width=True):
                st.session_state["nav"] = "register"
                st.rerun()

        with c3:
            if st.button("🔐 Forgot Password?", use_container_width=True):
                st.session_state["nav"] = "forgot_email"
                st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# REGISTER PAGE
# ══════════════════════════════════════════════════════════════════════════════
def page_register():
    _auth_topbar()
    _, mid, _ = st.columns([1, 1.4, 1])

    with mid:
        st.markdown("""
        <div class="auth-card" style="text-align:center;margin-top:1rem;">
          <div style="font-size:2rem">📝</div>
          <h2 style="margin:.35rem 0 .15rem 0;font-size:1.45rem">Create Account</h2>
          <p style="margin:0">Fill in your details to register.</p>
        </div>
        """, unsafe_allow_html=True)

        with st.form("reg_form", clear_on_submit=True):
            r_name = st.text_input("Full Name *", placeholder="Your full name")
            r_user = st.text_input("Username *", placeholder="Pick a username")
            r_email = st.text_input("Email *", placeholder="your@email.com")
            r_pw = st.text_input(
                "Password *",
                type="password",
                placeholder="Min 8 chars · 1 UPPER · 1 special"
            )
            r_pw2 = st.text_input("Confirm Password *", type="password")

            go = st.form_submit_button("Create Account", use_container_width=True)

            if go:
                if not r_name or not r_user or not r_email or not r_pw or not r_pw2:
                    st.error("All fields marked * are required.")

                elif r_pw != r_pw2:
                    st.error("Passwords do not match.")

                else:
                    user, err = _register(r_user, r_pw, r_name, r_email)
                    if err:
                        st.error(f"❌ {err}")
                    else:
                         # Auto login after registration
                        st.session_state["auth_user"] = user
                        st.session_state["session_id"] = uuid.uuid4().hex
                        st.session_state.pop("nav", None)
                        st.success("✅ Account created successfully! Redirecting to your dashboard...")
                        st.rerun()

        if st.button("← Back", use_container_width=True):
            st.session_state["nav"] = "landing"
            st.rerun()

# ══════════════════════════════════════════════════════════════════════════════
# FORGOT PASSWORD – STEP 1: enter email
# ══════════════════════════════════════════════════════════════════════════════
def page_forgot_email():
    _auth_topbar()
    _, mid, _ = st.columns([1, 1.4, 1])
    with mid:
        st.markdown("""
        <div style="background:#1a1a2e;border:1px solid #2a2a4e;border-radius:18px;
                    padding:2.2rem 2rem;margin-top:1.5rem">
          <div style="text-align:center;margin-bottom:1.5rem">
            <span style="font-size:2rem">📧</span>
            <h2 style="color:#e0e0ff;margin:.4rem 0 .2rem 0;font-size:1.4rem">Forgot Password</h2>
            <p style="color:#6060a0;font-size:.85rem;margin:0">
              Enter your registered email. We'll send a 6-digit OTP.
            </p>
          </div>
        </div>""", unsafe_allow_html=True)

        with st.form("fp_email_form"):
            email = st.text_input("Registered Email", placeholder="your@email.com")
            go    = st.form_submit_button("Send OTP  →", type="primary", use_container_width=True)
            if go:
                if not email:
                    st.error("Please enter your email.")
                else:
                    try:
                        SL = load_db()
                        from backend.database.models import User
                        db = SL()
                        u = db.query(User).filter(User.email == email, User.is_active == True).first()
                        db.close()
                        if not u:
                            st.error("No account found with that email.")
                        else:
                            otp = _generate_otp()
                            if not _store_otp(u.id, otp):
                                st.error("Could not generate OTP. Try again.")
                            else:
                                html = f"""
                                <div style="font-family:Arial;max-width:500px;margin:auto;
                                            border:1px solid #ddd;border-radius:10px;overflow:hidden">
                                  <div style="background:#1E3A8A;padding:20px;text-align:center">
                                    <h2 style="color:white;margin:0">🧠 DigitAI</h2>
                                  </div>
                                  <div style="padding:30px">
                                    <h3 style="color:#1E3A8A">Password Reset OTP</h3>
                                    <p>Hello <b>{u.full_name}</b>,</p>
                                    <p>Your one-time password (OTP) for resetting your DigitAI password is:</p>
                                    <div style="background:#EFF6FF;border:2px solid #2563EB;border-radius:10px;
                                                padding:20px;text-align:center;margin:20px 0">
                                      <span style="font-size:2.5rem;font-weight:800;color:#1D4ED8;
                                                   letter-spacing:8px">{otp}</span>
                                    </div>
                                    <p style="color:#666">This OTP is valid for <b>10 minutes</b>.</p>
                                    <p style="color:#999;font-size:12px">
                                      If you did not request this, please ignore this email.
                                    </p>
                                  </div>
                                </div>"""
                                smtp_err = _send_email(email, "DigitAI – Password Reset OTP", html)
                                if smtp_err:
                                    st.warning(f"⚠️ Email could not be sent ({smtp_err}). "
                                               f"For testing, your OTP is: **{otp}**")
                                else:
                                    st.success("✅ OTP sent to your email!")
                                st.session_state["fp_user_id"] = u.id
                                st.session_state["fp_email"]   = email
                                st.session_state["nav"] = "forgot_otp"
                                st.rerun()
                    except Exception as e:
                        st.error(f"Error: {e}")

        if st.button("← Back to Login"):
            st.session_state["nav"] = "login"; st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# FORGOT PASSWORD – STEP 2: enter OTP
# ══════════════════════════════════════════════════════════════════════════════
def page_forgot_otp():
    _auth_topbar()
    _, mid, _ = st.columns([1, 1.4, 1])
    with mid:
        email = st.session_state.get("fp_email", "your email")
        st.markdown(f"""
        <div style="background:#1a1a2e;border:1px solid #2a2a4e;border-radius:18px;
                    padding:2.2rem 2rem;margin-top:1.5rem">
          <div style="text-align:center;margin-bottom:1.5rem">
            <span style="font-size:2rem">🔢</span>
            <h2 style="color:#e0e0ff;margin:.4rem 0 .2rem 0;font-size:1.4rem">Enter OTP</h2>
            <p style="color:#6060a0;font-size:.85rem;margin:0">
              We sent a 6-digit OTP to <b style="color:#4a9eff">{email}</b>.<br>
              Valid for 10 minutes.
            </p>
          </div>
        </div>""", unsafe_allow_html=True)

        with st.form("fp_otp_form"):
            otp = st.text_input("Enter 6-digit OTP", placeholder="e.g. 482931", max_chars=6)
            go  = st.form_submit_button("Verify OTP  →", type="primary", use_container_width=True)
            if go:
                uid = st.session_state.get("fp_user_id")
                if not otp or len(otp) != 6:
                    st.error("Enter a valid 6-digit OTP.")
                elif not uid:
                    st.error("Session expired. Start again."); st.session_state["nav"] = "forgot_email"; st.rerun()
                else:
                    err = _verify_otp(uid, otp)
                    if err:
                        st.error(f"❌ {err}")
                    else:
                        st.success("✅ OTP verified!")
                        st.session_state["nav"] = "forgot_reset"; st.rerun()

        c1, c2 = st.columns(2)
        with c1:
            if st.button("← Back"):
                st.session_state["nav"] = "forgot_email"; st.rerun()
        with c2:
            if st.button("Resend OTP"):
                st.session_state["nav"] = "forgot_email"; st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# FORGOT PASSWORD – STEP 3: set new password
# ══════════════════════════════════════════════════════════════════════════════
def page_forgot_reset():
    _auth_topbar()
    _, mid, _ = st.columns([1, 1.4, 1])
    with mid:
        st.markdown("""
        <div style="background:#1a1a2e;border:1px solid #2a2a4e;border-radius:18px;
                    padding:2.2rem 2rem;margin-top:1.5rem">
          <div style="text-align:center;margin-bottom:1.5rem">
            <span style="font-size:2rem">🔒</span>
            <h2 style="color:#e0e0ff;margin:.4rem 0 .2rem 0;font-size:1.4rem">Set New Password</h2>
            <p style="color:#6060a0;font-size:.85rem;margin:0">
              OTP verified. Choose a strong new password.
            </p>
          </div>
        </div>""", unsafe_allow_html=True)

        with st.form("fp_reset_form"):
            np1 = st.text_input("New Password", type="password",
                                placeholder="Min 8 chars · 1 UPPER · 1 special")
            np2 = st.text_input("Confirm New Password", type="password")
            go  = st.form_submit_button("Reset Password  →", type="primary", use_container_width=True)
            if go:
                uid = st.session_state.get("fp_user_id")
                if not np1 or not np2:
                    st.error("Both fields required.")
                elif np1 != np2:
                    st.error("Passwords do not match.")
                else:
                    err = _validate_password(np1)
                    if err:
                        st.error(f"❌ {err}")
                    elif not uid:
                        st.error("Session expired."); st.session_state["nav"] = "forgot_email"; st.rerun()
                    else:
                        email = _reset_password(uid, np1)
                        # Send success email
                        if email:
                            html = """
                            <div style="font-family:Arial;max-width:500px;margin:auto;
                                        border:1px solid #ddd;border-radius:10px;overflow:hidden">
                              <div style="background:#166534;padding:20px;text-align:center">
                                <h2 style="color:white;margin:0">🧠 DigitAI</h2>
                              </div>
                              <div style="padding:30px">
                                <h3 style="color:#166534">✅ Password Reset Successful</h3>
                                <p>Your DigitAI password has been successfully reset.</p>
                                <p>You can now log in with your new password.</p>
                                <p style="color:#999;font-size:12px;margin-top:20px">
                                  If you did not perform this action, contact support immediately.
                                </p>
                              </div>
                            </div>"""
                            _send_email(email, "DigitAI – Password Reset Successful", html)
                        for k in ["fp_user_id", "fp_email"]:
                            st.session_state.pop(k, None)
                        st.success("✅ Password reset successfully! Redirecting to login…")
                        _dt.time  # just to force a moment of delay perception
                        st.session_state["nav"] = "login"; st.rerun()



def _render_dark_mode_toggle(widget_key: str):
    """Render dark-mode toggle with a unique widget key and sync it to st.session_state['dark_mode']."""
    current = bool(st.session_state.get("dark_mode", True))
    new_value = st.toggle("🌙 Dark Mode", value=current, key=widget_key)
    if bool(new_value) != current:
        st.session_state["dark_mode"] = bool(new_value)
        st.rerun()


def _auth_topbar():
    mode = "Dark" if st.session_state.get("dark_mode", True) else "Light"
    st.markdown(f"""
    <div class="app-topbar">
      <div class="brand-wrap">
        <div class="brand-logo">🧠</div>
        <div>
          <div class="brand-title">DigitAI</div>
          <div class="brand-subtitle">Advanced Digit Recognition • {mode} Mode</div>
        </div>
      </div>
      <div class="topbar-user">CNN • OCR • Documents</div>
    </div>
    """, unsafe_allow_html=True)
    c1, c2 = st.columns([6, 1])
    with c2:
        _render_dark_mode_toggle("dark_mode_toggle_auto")


def _render_navbar(user):
    role_html = ('<span class="role-admin">ADMIN</span>' if user["role"] == "admin"
                 else '<span class="role-client">CLIENT</span>')
    mode = "Dark" if st.session_state.get("dark_mode", True) else "Light"
    st.markdown(f"""
    <div class="app-topbar">
      <div class="brand-wrap">
        <div class="brand-logo">🧠</div>
        <div>
          <div class="brand-title">DigitAI</div>
          <div class="brand-subtitle">Canvas • Upload • OCR • Analytics • {mode} Mode</div>
        </div>
      </div>
      <div style="display:flex;align-items:center;gap:.75rem;">
        {role_html}
        <span class="topbar-user">👤 {user["full_name"]}</span>
      </div>
    </div>
    """, unsafe_allow_html=True)

    nav_c1, nav_c2, nav_c3 = st.columns([7, 1.25, 1])
    with nav_c2:
        _render_dark_mode_toggle("dark_mode_toggle_nav")
    with nav_c3:
        if st.button("Logout", key="logout_btn", use_container_width=True):
            st.session_state["logout_confirm"] = True
            st.rerun()

    if st.session_state.get("logout_confirm"):
        st.warning("⚠️ Are you sure you want to logout?")
        y_col, n_col, _ = st.columns([1.2, 1.2, 7.6])
        with y_col:
            if st.button("✅ Yes, Logout", type="primary", use_container_width=True):
                st.session_state.clear()
                st.rerun()
        with n_col:
            if st.button("❌ Cancel", use_container_width=True):
                st.session_state["logout_confirm"] = False
                st.rerun()
        st.stop()


def page_predict(user):
    st.markdown("## ✏️ Predict a Digit")
    st.markdown("Draw a digit on the canvas **or** upload an image, then click **Predict**.")
    col_left, col_right = st.columns([1, 1], gap="large")

    with col_left:
        input_tab, upload_tab = st.tabs(["🖊️ Draw", "📁 Upload"])
        with input_tab:
            st.markdown('<div class="section-title">Draw a digit (0–9)</div>', unsafe_allow_html=True)
            brush = st.slider("Brush size", 8, 30, 16, label_visibility="collapsed")
            canvas_result = st_canvas(
                fill_color="rgba(255,255,255,0)", stroke_width=brush,
                stroke_color="#000000", background_color="#FFFFFF",
                height=280, width=280, drawing_mode="freedraw", key="canvas",
            )
            input_image = None
            if canvas_result.image_data is not None:
                arr = canvas_result.image_data.astype(np.uint8)
                if arr.max() > 0: input_image = arr
            input_type = "canvas"

        with upload_tab:
            st.markdown('<div class="section-title">Upload an image</div>', unsafe_allow_html=True)
            uploaded = st.file_uploader("Choose image",
                                        type=["png","jpg","jpeg","bmp","webp"],
                                        label_visibility="collapsed")
            if uploaded:
                pil_img = ImageOps.exif_transpose(Image.open(uploaded)).convert("RGB")
                st.image(pil_img, caption="Uploaded", width=280)
                input_image = np.array(pil_img)[:, :, ::-1]
                input_type  = "upload"

        st.markdown('<div class="section-title">Model & Options</div>', unsafe_allow_html=True)
        sel_model = st.selectbox("Model", AVAILABLE_MODELS,
                                 format_func=lambda m: MODEL_LABELS.get(m, m),
                                 index=min(1, len(AVAILABLE_MODELS)-1),
                                 label_visibility="collapsed")
        xai_on    = st.checkbox("Show XAI heatmaps (Grad-CAM & LIME)", value=True)
        pred_btn  = st.button("🔮  Predict Digit", type="primary",
                              use_container_width=True, disabled=(input_image is None))

    with col_right:
        if pred_btn and input_image is not None:
            with st.spinner("Running inference…"):
                try:
                    pfn, gcfn, limefn, fgsmfn, reg = load_prediction_service()
                    result = pfn(input_image, sel_model)
                    st.session_state.update({
                        "last_result": result, "last_image": input_image,
                        "last_model": sel_model, "last_input_type": input_type, "xai_on": xai_on,
                    })
                    if xai_on:
                        st.session_state["gradcam_path"] = gcfn(input_image, sel_model)
                        st.session_state["lime_path"]    = limefn(input_image, sel_model)
                    _save_prediction(result, input_image, input_type, xai_on, user)
                except Exception as e:
                    st.error(f"Prediction failed: {e}")

        result = st.session_state.get("last_result")
        if result: _render_results(result)
        else:      st.info("👈 Draw or upload a digit, then click **Predict Digit**.")

    st.markdown("---")
    st.markdown("## 📝 OCR Text Prediction")
    st.markdown("Draw or upload text / alphabet / word, then click **Predict OCR Text**.")

    # ── OCR model selector ────────────────────────────────────────────────────
    from backend.services.ocr_prediction_service import (
        get_available_ocr_models, load_ocr_model_from_path, predict_full_image_ocr
    )
    _ocr_model_map = get_available_ocr_models()
    _ocr_labels    = list(_ocr_model_map.keys())

    ocr_left, ocr_right = st.columns([1, 1], gap="large")

    with ocr_left:
        ocr_draw_tab, ocr_upload_tab = st.tabs(["🖊️ Draw", "📁 Upload"])

        ocr_input_image = None

        with ocr_draw_tab:
            st.markdown('<div class="section-title">Draw a letter / word</div>', unsafe_allow_html=True)
            ocr_brush = st.slider("Brush size", 4, 24, 10,
                                  key="ocr_brush", label_visibility="collapsed")
            ocr_canvas = st_canvas(
                fill_color="rgba(255,255,255,0)",
                stroke_width=ocr_brush,
                stroke_color="#000000",
                background_color="#FFFFFF",
                height=120,
                width=560,
                drawing_mode="freedraw",
                key="ocr_canvas",
            )
            if ocr_canvas.image_data is not None:
                _arr = ocr_canvas.image_data.astype(np.uint8)
                if _arr.max() > 0:
                    ocr_input_image = Image.fromarray(_arr)
            ocr_input_type = "canvas"

        with ocr_upload_tab:
            st.markdown('<div class="section-title">Upload an image</div>', unsafe_allow_html=True)
            uploaded_ocr = st.file_uploader(
                "Upload alphabet / word / OCR image",
                type=["png", "jpg", "jpeg"],
                key="ocr_upload_predict",
                label_visibility="collapsed",
            )
            if uploaded_ocr is not None:
                ocr_input_image = ImageOps.exif_transpose(Image.open(uploaded_ocr))
                st.image(ocr_input_image, caption="Uploaded", use_container_width=True)
                ocr_input_type = "upload"

        # OCR model selector
        st.markdown('<div class="section-title">OCR Model</div>', unsafe_allow_html=True)
        if _ocr_labels:
            _sel_ocr_label = st.selectbox(
                "OCR Model",
                _ocr_labels,
                index=0,
                label_visibility="collapsed",
                key="ocr_model_select_predict",
            )
            _ocr_model_path = str(_ocr_model_map[_sel_ocr_label])
        else:
            _ocr_model_path = None
            st.warning("No OCR model found in models/. Please add ocr_prediction_model.keras.")

        st.markdown('<div class="section-title">OCR Reading Mode</div>', unsafe_allow_html=True)
        _ocr_mode_label = st.selectbox(
            "OCR Reading Mode",
            ["Printed Text / Screenshot", "Words or Line", "Auto Detect", "Separate Characters", "Full Image"],
            index=1,
            label_visibility="collapsed",
            key="ocr_mode_select_predict",
            help="Choose how the app should read your image. Use Auto Detect if you are not sure.",
        )
        _ocr_mode_map = {
            "Printed Text / Screenshot": "printed_paragraph",
            "Words or Line": "line",
            "Auto Detect": "auto",
            "Separate Characters": "character_grid",
            "Full Image": "whole",
        }
        _ocr_force_upper = st.checkbox(
            "Convert result to UPPERCASE",
            value=False,
            key="ocr_uppercase_predict",
            help="Useful when your image contains only capital letters."
        )

        ocr_pred_btn = st.button(
            "📝  Predict OCR Text", type="primary",
            use_container_width=True,
            disabled=(ocr_input_image is None),
            key="ocr_predict_btn_predict",
        )

    with ocr_right:
        if ocr_pred_btn and ocr_input_image is not None:
            with st.spinner("Running OCR…"):
                try:
                    _ocr_mode = _ocr_mode_map.get(_ocr_mode_label, "line")
                    # Printed Paragraph and Handwritten Paragraph modes load their own engine/model inside the OCR service.
                    _external_ocr_modes = {"printed_paragraph"}
                    _ocr_loaded = None if _ocr_mode in _external_ocr_modes else (load_ocr_model_from_path(_ocr_model_path) if _ocr_model_path else None)
                    ocr_result = predict_full_image_ocr(
                        ocr_input_image,
                        model=_ocr_loaded,
                        mode=_ocr_mode,
                    )
                    if _ocr_force_upper:
                        ocr_result["text"] = str(ocr_result.get("text", "")).upper()
                        for _line in ocr_result.get("lines", []):
                            if "text" in _line:
                                _line["text"] = str(_line.get("text", "")).upper()
                            for _ch in _line.get("char_details", []):
                                if "char" in _ch:
                                    _ch["char"] = str(_ch.get("char", "")).upper()
                    st.session_state["last_ocr_result"] = ocr_result
                except Exception as e:
                    st.error(f"OCR failed: {e}")

        ocr_result = st.session_state.get("last_ocr_result")
        if ocr_result:
            st.markdown('<div class="section-title">OCR Result</div>', unsafe_allow_html=True)
            predicted_text = ocr_result.get("text", "")
            confidence     = ocr_result.get("confidence", 0.0) or 0.0
            # Safety: custom CRNN returns 0..1, but some external OCR engines may return 0..100.
            # Normalize before displaying so confidence never becomes 7000%.
            try:
                confidence = float(confidence)
                if confidence > 1.0:
                    confidence = confidence / 100.0
            except Exception:
                confidence = 0.0
            pill_cls = "high" if confidence >= 0.85 else "mid" if confidence >= 0.60 else "low"

            st.success("Prediction:")
            st.code(predicted_text if predicted_text else "(no text detected)")
            if ocr_result.get("segmentation_mode"):
                st.caption(f"OCR mode used: {ocr_result.get('segmentation_mode')}")

            if ocr_result.get("confidence") is not None:
                st.markdown(f"""
                <div class="conf-pill {pill_cls}" style="display:inline-block;margin:.4rem 0">
                    {confidence:.2%} confidence
                </div>""", unsafe_allow_html=True)

            with st.expander("Line-wise & Alphabet-wise details", expanded=True):
                lines = ocr_result.get("lines", [])
                if lines:
                    for line in lines:
                        lconf = line.get("confidence", 0) or 0
                        try:
                            lconf = float(lconf)
                            if lconf > 1.0:
                                lconf = lconf / 100.0
                        except Exception:
                            lconf = 0.0
                        st.markdown(f"**Line {line.get('line_index', 0) + 1}**: `{line.get('text', '')}` ({lconf:.2%} avg)")

                        char_details = line.get("char_details", [])
                        if char_details:
                            import pandas as pd
                            df_chars = pd.DataFrame(char_details)
                            df_chars["confidence"] = pd.to_numeric(df_chars["confidence"], errors="coerce").fillna(0.0)
                            df_chars.loc[df_chars["confidence"] > 1.0, "confidence"] = df_chars.loc[df_chars["confidence"] > 1.0, "confidence"] / 100.0
                            df_chars["Confidence (%)"] = (df_chars["confidence"] * 100).round(2)
                            df_chars.rename(columns={"char": "Alphabet"}, inplace=True)

                            def _style_char_conf(val):
                                try:
                                    v = float(val)
                                    if v >= 85: return "background-color:#0d3a1a; color:#4ade80"
                                    if v >= 60: return "background-color:#3a2e00; color:#fbbf24"
                                    return "background-color:#3a0d0d; color:#f87171"
                                except: return ""

                            st.dataframe(
                                df_chars[["Alphabet", "Confidence (%)"]].style.applymap(
                                    _style_char_conf, subset=["Confidence (%)"]
                                ),
                                use_container_width=True,
                                hide_index=True
                            )
                else:
                    st.write("No line segmentation data available.")
        else:
            st.info("👈 Draw or upload text, then click **Predict OCR Text**.")





def _confidence_to_percent(confidence):
    """Convert model confidence to percentage. Accepts 0-1 or 0-100 values."""
    try:
        value = float(confidence)
        if value <= 1.0:
            value *= 100.0
        return max(0.0, min(value, 100.0))
    except Exception:
        return None


def _render_result_and_accuracy(result_text, confidence=None, result_label="Result"):
    """
    Compact output renderer.
    Shows only prediction result and accuracy/confidence as requested.
    """
    result_text = "" if result_text is None else str(result_text)
    percent = _confidence_to_percent(confidence)

    st.markdown('<div class="section-title">Prediction Output</div>', unsafe_allow_html=True)

    c1, c2 = st.columns([2, 1])

    with c1:
        st.markdown(f"**{result_label}**")
        if len(result_text) > 180:
            st.text_area(
                "Prediction Result",
                result_text if result_text.strip() else "(no result)",
                height=180,
                label_visibility="collapsed",
                disabled=True,
            )
        else:
            st.code(result_text if result_text.strip() else "(no result)")

    with c2:
        if percent is None:
            st.metric("Accuracy", "N/A")
        else:
            st.metric("Accuracy", f"{percent:.2f}%")


def _render_results(result):
    """
    Compact single-digit result.
    Only shows result and accuracy/confidence.
    """
    digit = result.get("predicted_digit", "")
    conf = result.get("confidence", None)
    _render_result_and_accuracy(digit, conf, result_label="Predicted Digit")



def _save_prediction(result, image_arr, input_type, xai_on, user):
    """
    Save single digit prediction in the old PredictionLog table
    and also in UniversalPredictionLog so user/admin history can show all records together.
    """
    try:
        SL = load_db()
        from backend.database.models import PredictionLog
        import cv2

        db = SL()

        ud = ROOT / "frontend" / "static" / "uploads"
        ud.mkdir(parents=True, exist_ok=True)

        fname = f"upload_{uuid.uuid4().hex[:10]}.png"

        img_bgr = (
            image_arr
            if getattr(image_arr, "ndim", 0) == 3
            else cv2.cvtColor(image_arr, cv2.COLOR_GRAY2BGR)
        )

        cv2.imwrite(str(ud / fname), img_bgr)

        db.add(PredictionLog(
            session_id=st.session_state.get("session_id", uuid.uuid4().hex),
            user_id=user["id"],
            username=user["username"],
            input_type=input_type,
            image_path=f"static/uploads/{fname}",
            model_used=result.get("model_used"),
            predicted_digit=result.get("predicted_digit"),
            confidence=result.get("confidence"),
            top3_predictions=result.get("top3_predictions"),
            all_probabilities=result.get("all_probabilities"),
            gradcam_path=st.session_state.get("gradcam_path") if xai_on else None,
            processing_time_ms=result.get("processing_time_ms"),
        ))

        db.commit()

        latest = db.query(PredictionLog).order_by(PredictionLog.id.desc()).first()
        if latest:
            st.session_state["last_log_id"] = latest.id

        db.close()

        # Also save in universal history table.
        try:
            _save_universal_prediction(
                category="digit",
                user=user,
                input_source=input_type,
                model_used=result.get("model_used", ""),
                predicted_digit=result.get("predicted_digit"),
                confidence=result.get("confidence"),
                top3_predictions=result.get("top3_predictions"),
                image_arr=image_arr,
                processing_time_ms=result.get("processing_time_ms", 0),
                extra_data={
                    "all_probabilities": result.get("all_probabilities"),
                    "saved_from": "single_digit_prediction",
                },
            )
        except Exception:
            pass

    except Exception as e:
        st.warning(f"Could not save prediction history: {e}")


def _submit_feedback(is_correct, true_label):
    try:
        SL = load_db()
        from backend.database.models import PredictionLog
        lid = st.session_state.get("last_log_id")
        if not lid: return
        db = SL()
        r = db.query(PredictionLog).filter(PredictionLog.id == lid).first()
        if r: r.is_correct = is_correct; r.true_label = true_label; db.commit()
        db.close()
    except: pass


# ══════════════════════════════════════════════════════════════════════════════
# UNIVERSAL HISTORY SAVE  (digit / multidigit / alphabet / script)
# ══════════════════════════════════════════════════════════════════════════════
def _save_universal_prediction(
    category: str,          # 'digit' | 'multidigit' | 'alphabet' | 'script'
    user: dict,
    input_source: str = "canvas",   # 'canvas' | 'upload' | 'batch' | 'document'
    model_used: str = "",
    ocr_mode: str = "",
    # digit fields
    predicted_digit=None,
    confidence=None,
    top3_predictions=None,
    # multidigit fields
    multidigit_result: str = "",
    digit_count: int = 0,
    # ocr/alphabet/script fields
    ocr_text: str = "",
    ocr_confidence: float = 0.0,
    line_count: int = 0,
    # image / doc
    image_arr=None,
    source_filename: str = "",
    # timing
    processing_time_ms: float = 0.0,
    # extra
    extra_data=None,
):
    """Save any prediction type to universal_prediction_logs."""
    try:
        SL = load_db()
        from backend.database.models import UniversalPredictionLog
        import cv2
        db = SL()

        image_path = None
        if image_arr is not None:
            try:
                ud = ROOT / "frontend" / "static" / "uploads"
                ud.mkdir(parents=True, exist_ok=True)
                fname = f"univ_{uuid.uuid4().hex[:10]}.png"
                img_bgr = (
                    image_arr if image_arr.ndim == 3
                    else cv2.cvtColor(image_arr, cv2.COLOR_GRAY2BGR)
                )
                cv2.imwrite(str(ud / fname), img_bgr)
                image_path = f"static/uploads/{fname}"
            except Exception:
                pass

        db.add(UniversalPredictionLog(
            session_id=st.session_state.get("session_id", uuid.uuid4().hex),
            user_id=user["id"],
            username=user["username"],
            prediction_category=category,
            input_source=input_source,
            source_filename=source_filename or None,
            image_path=image_path,
            model_used=model_used or None,
            ocr_mode=ocr_mode or None,
            predicted_digit=predicted_digit,
            confidence=float(confidence) if confidence is not None else None,
            top3_predictions=top3_predictions,
            multidigit_result=multidigit_result or None,
            digit_count=int(digit_count) if digit_count else None,
            ocr_text=ocr_text or None,
            ocr_confidence=float(ocr_confidence) if ocr_confidence else None,
            line_count=int(line_count) if line_count else None,
            processing_time_ms=float(processing_time_ms) if processing_time_ms else None,
            extra_data=extra_data,
        ))
        db.commit()
        db.close()
    except Exception:
        pass




# ══════════════════════════════════════════════════════════════════════════════
# PAGE: BATCH PREDICTION  ← NEW
# ══════════════════════════════════════════════════════════════════════════════

def page_batch_predict(user):
    st.markdown("## 📦 Batch Prediction")
    st.markdown("Upload multiple single-digit images. Output will show only result and accuracy for each image.")

    b_model = st.selectbox(
        "Model for batch",
        AVAILABLE_MODELS,
        format_func=lambda m: MODEL_LABELS.get(m, m),
        index=min(1, len(AVAILABLE_MODELS) - 1),
        key="batch_model",
    )

    uploaded_files = st.file_uploader(
        "Select images",
        type=["png", "jpg", "jpeg", "bmp", "tiff", "tif", "webp"],
        accept_multiple_files=True,
        help="Upload one or more handwritten digit images.",
        key="batch_digit_uploads",
    )

    if not uploaded_files:
        st.info("Upload one or more digit images to begin batch prediction.")
        return

    run_btn = st.button("🚀 Run Batch Prediction", type="primary", use_container_width=True)

    if not run_btn and "batch_results" not in st.session_state:
        return

    if run_btn:
        pfn, *_ = load_prediction_service()
        results = []
        progress = st.progress(0, text="Processing images...")

        for idx, f in enumerate(uploaded_files):
            progress.progress(idx / max(len(uploaded_files), 1), text=f"Processing {f.name}...")

            try:
                f.seek(0)
                raw = f.read()
                img_pil = ImageOps.exif_transpose(Image.open(io.BytesIO(raw))).convert("RGB")
                img_arr = np.array(img_pil)[:, :, ::-1]

                result = pfn(img_arr, b_model)
                accuracy = _confidence_to_percent(result.get("confidence"))

                results.append({
                    "Image": f.name,
                    "Result": result.get("predicted_digit"),
                    "Accuracy (%)": round(accuracy or 0, 2),
                })

                _save_prediction(result, img_arr, "batch_upload", False, user)

            except Exception as e:
                results.append({
                    "Image": f.name,
                    "Result": "Error",
                    "Accuracy (%)": 0.0,
                })

        progress.progress(1.0, text="Batch complete.")
        st.session_state["batch_results"] = results

    results = st.session_state.get("batch_results", [])

    if results:
        st.markdown("### Batch Output")
        st.dataframe(
            pd.DataFrame(results)[["Image", "Result", "Accuracy (%)"]],
            use_container_width=True,
            hide_index=True,
            height=360,
        )

        stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
        st.download_button(
            "⬇️ Download CSV",
            pd.DataFrame(results).to_csv(index=False),
            f"batch_results_{stamp}.csv",
            "text/csv",
            use_container_width=True,
        )

    if st.button("🔄 Clear Batch Results", use_container_width=True):
        st.session_state.pop("batch_results", None)
        st.rerun()


def _batch_pdf(df: pd.DataFrame, model: str, avg_conf: float, total_ms: float) -> bytes:
    from io import BytesIO
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    buf  = BytesIO()
    doc  = SimpleDocTemplate(buf, pagesize=landscape(A4),
                             leftMargin=1*cm, rightMargin=1*cm,
                             topMargin=1.5*cm, bottomMargin=1.5*cm)
    styl = getSampleStyleSheet()
    now  = _dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    ok   = len(df[df["Status"] == "OK"]) if "Status" in df.columns else len(df)
    elems = [
        Paragraph("<b>DigitAI – Batch Prediction Report</b>", styl["Title"]),
        Paragraph(
            f"Generated: {now}  |  Model: {model}  |  Images: {len(df)}  |  "
            f"Success: {ok}  |  Avg Confidence: {avg_conf:.1f}%  |  Total Time: {total_ms:.0f} ms",
            styl["Normal"],
        ),
        Spacer(1, .4*cm),
    ]
    data = [list(df.columns)] + df.astype(str).values.tolist()
    tbl  = Table(data, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND",    (0,0),(-1,0),  colors.HexColor("#1E3A8A")),
        ("TEXTCOLOR",     (0,0),(-1,0),  colors.white),
        ("FONTNAME",      (0,0),(-1,0),  "Helvetica-Bold"),
        ("FONTSIZE",      (0,0),(-1,-1), 7.5),
        ("ROWBACKGROUNDS",(0,1),(-1,-1), [colors.white, colors.HexColor("#EFF6FF")]),
        ("GRID",          (0,0),(-1,-1), .4, colors.HexColor("#CBD5E1")),
        ("ALIGN",         (0,0),(-1,-1), "CENTER"),
        ("VALIGN",        (0,0),(-1,-1), "MIDDLE"),
        ("TOPPADDING",    (0,0),(-1,-1), 4),
        ("BOTTOMPADDING", (0,0),(-1,-1), 4),
    ]))
    elems.append(tbl)
    doc.build(elems)
    return buf.getvalue()


def analyze_image_quality(img):
    import cv2

    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    brightness = np.mean(gray)
    blur_score = cv2.Laplacian(gray, cv2.CV_64F).var()
    return {
        "brightness": brightness,
        "blur_score": blur_score,
        "is_dark": brightness < 70,
        "is_bright": brightness > 200,
        "is_blurry": blur_score < 100,
    }


def _nms_boxes(boxes, overlap_thresh=0.3):
    """Remove heavily overlapping / nested bounding boxes (IoU-based NMS)."""
    if not boxes:
        return []
    boxes = sorted(boxes, key=lambda b: b[2] * b[3], reverse=True)  # largest first
    keep = []
    for box in boxes:
        x, y, w, h = box
        suppressed = False
        for kx, ky, kw, kh in keep:
            # Intersection
            ix1 = max(x, kx);  iy1 = max(y, ky)
            ix2 = min(x+w, kx+kw);  iy2 = min(y+h, ky+kh)
            iw = max(0, ix2-ix1);  ih = max(0, iy2-iy1)
            inter = iw * ih
            union = w*h + kw*kh - inter
            if union > 0 and inter / union > overlap_thresh:
                suppressed = True
                break
        if not suppressed:
            keep.append(box)
    return keep



def page_multi_digit_predict(user):
    st.markdown("## 🔢 Universal Digit Detection")
    st.info(
        "Upload any digit image. This page uses common detection for black, white, colored, pencil, script, "
        "badge/icon, dark-background, and multi-row digit images."
    )

    # Select best available model by default.
    default_index = 0
    if "cnn_deep" in AVAILABLE_MODELS:
        default_index = AVAILABLE_MODELS.index("cnn_deep")
    elif "cnn_medium" in AVAILABLE_MODELS:
        default_index = AVAILABLE_MODELS.index("cnn_medium")

    selected_model = st.selectbox(
        "Select Model",
        AVAILABLE_MODELS,
        format_func=lambda m: MODEL_LABELS.get(m, m),
        index=default_index,
        key="universal_digit_model",
    )

    c1, c2, c3 = st.columns(3)
    with c1:
        confidence_threshold = st.slider(
            "Confidence Threshold",
            min_value=0,
            max_value=100,
            value=20,
            step=5,
            help="Lower value keeps more predicted digits. Higher value marks doubtful digits as low confidence.",
            key="universal_digit_confidence",
        )
    with c2:
        max_digits = st.slider(
            "Maximum Digits",
            min_value=1,
            max_value=200,
            value=100,
            step=1,
            key="universal_digit_max_digits",
        )
    with c3:
        show_debug = st.checkbox(
            "Show mask and crops",
            value=True,
            key="universal_digit_debug",
        )

    uploaded = st.file_uploader(
        "Upload digit image",
        type=["png", "jpg", "jpeg", "bmp", "webp"],
        key="universal_digit_upload",
    )

    if uploaded is None:
        st.info("Upload an image containing digits to start detection.")
        return

    pil_img = ImageOps.exif_transpose(Image.open(uploaded)).convert("RGB")
    img_rgb = np.array(pil_img)
    st.image(pil_img, caption="Uploaded Image", use_container_width=True)

    if not st.button("🔍 Detect Digits", type="primary", key="universal_digit_detect_btn"):
        return

    try:
        from preprocessing.universal_digit_detector import detect_digits_universal
        from backend.services.prediction_service import registry

        model = registry.get(selected_model)
        if model is None:
            st.error(f"Model not found: {selected_model}. Please check your models folder.")
            return

        with st.spinner("Detecting and predicting digits..."):
            result = detect_digits_universal(
                img_rgb,
                model=model,
                confidence_threshold=float(confidence_threshold),
                max_digits=int(max_digits),
            )

        if not result.get("success"):
            st.error(result.get("message", "No digits detected."))
            if show_debug and result.get("mask_image") is not None:
                st.image(result["mask_image"], caption="Combined Mask", use_container_width=True, clamp=True)
            return

        st.success(f"✅ Detected digit-like items: **{result['digit_count']}**")
        
        # ── Save to universal history ──────────────────────
        _save_universal_prediction(
            category="multidigit",
            user=user,
            input_source="upload",
            model_used=selected_model,
            multidigit_result=str(result.get("number_prediction") or result.get("prediction", "")),
            digit_count=result.get("digit_count", 0),
            image_arr=img_rgb,
            extra_data={"rows": result.get("number_rows", [])},
        )


        st.markdown("### Final Output")
        # Number-wise output keeps multi-digit numbers like 10, 11, 123 as complete numbers.
        final_number_output = result.get("number_prediction") or result.get("prediction", "")
        st.code(final_number_output)

        digit_sequence_output = result.get("prediction", "")
        if digit_sequence_output and digit_sequence_output != final_number_output:
            with st.expander("Digit sequence output without number grouping"):
                st.code(digit_sequence_output)

        raw_number_prediction = result.get("raw_number_prediction", "")
        if raw_number_prediction and raw_number_prediction != final_number_output:
            with st.expander("Raw model output including low-confidence digits"):
                st.code(raw_number_prediction)

        st.markdown("### Detected Boxes")
        st.image(result["annotated_image"], caption="Detected Digits", use_container_width=True)

        # Row-wise number output
        number_rows = result.get("number_rows", [])
        if number_rows:
            st.markdown("### Row-wise Numbers")
            for idx, row in enumerate(number_rows, start=1):
                row_text = "   ".join(str(item) for item in row)
                st.write(f"Row {idx}: `{row_text}`")

        # Also show individual digits if needed.
        rows = result.get("rows", [])
        if rows:
            with st.expander("Individual detected digits"):
                for idx, row in enumerate(rows, start=1):
                    row_text = " ".join(str(item.get("digit", "?")) for item in row)
                    st.write(f"Row {idx}: `{row_text}`")

        # Details table
        table_rows = []
        for d in result.get("digits", []):
            top3 = d.get("top3_predictions", [])
            table_rows.append({
                "Position": d.get("position"),
                "Digit": d.get("digit"),
                "Confidence (%)": d.get("confidence"),
                "Status": d.get("status"),
                "Mask Used": d.get("mask_name"),
                "Box": str(d.get("box")),
                "Top-3": " | ".join(f"{p.get('digit')}:{p.get('confidence')}%" for p in top3),
            })
        if table_rows:
            st.markdown("### Prediction Details")
            st.dataframe(pd.DataFrame(table_rows), use_container_width=True, hide_index=True)

        if show_debug:
            st.markdown("### Debug / Preprocessing")
            d1, d2 = st.columns(2)
            with d1:
                st.image(result["mask_image"], caption="Combined Digit Mask", use_container_width=True, clamp=True)
            with d2:
                st.caption("If boxes are correct but digit values are wrong, improve/retrain the CNN model with similar images.")

            crops = [d for d in result.get("digits", []) if d.get("canvas28") is not None]
            if crops:
                st.markdown("### 28×28 Crops Sent to CNN")
                cols = st.columns(min(10, len(crops)))
                for i, d in enumerate(crops):
                    with cols[i % len(cols)]:
                        st.image(
                            d["canvas28"],
                            caption=f"{d.get('digit')} ({float(d.get('confidence', 0)):.1f}%)",
                            width=85,
                            clamp=True,
                        )

        st.warning(
            "Important: this detector finds digit-like regions. If the image also contains letters/symbols, "
            "a digit-only CNN may still classify them as digits. For true 'every type' accuracy, retrain with "
            "colored/script/font digits and add a non-digit rejection class."
        )

    except ModuleNotFoundError as e:
        st.error("Missing detector file. Make sure preprocessing/universal_digit_detector.py exists.")
        st.code(str(e))
    except Exception as e:
        st.error(f"Universal digit detection failed: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: CLIENT DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════════════════════
# DOCUMENT INPUT HELPERS FOR OCR PAGE
# ══════════════════════════════════════════════════════════════════════════════

def _extract_pdf_text_and_images(uploaded_file, max_pages=5, zoom=2.0):
    """Return direct PDF text and rendered page images for OCR."""
    try:
        import fitz  # PyMuPDF
    except ModuleNotFoundError:
        raise RuntimeError("PDF support needs PyMuPDF. Install it with: pip install pymupdf")

    data = uploaded_file.getvalue()
    doc = fitz.open(stream=data, filetype="pdf")
    text_parts = []
    page_images = []

    total_pages = min(int(max_pages), doc.page_count)
    matrix = fitz.Matrix(float(zoom), float(zoom))

    for page_index in range(total_pages):
        page = doc.load_page(page_index)

        page_text = page.get_text("text").strip()
        if page_text:
            text_parts.append(f"--- PDF Page {page_index + 1} ---\n{page_text}")

        pix = page.get_pixmap(matrix=matrix, alpha=False)
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        page_images.append((f"PDF page {page_index + 1}", img))

    doc.close()
    return "\n\n".join(text_parts).strip(), page_images


def _extract_docx_text_and_images(uploaded_file, max_images=25):
    """Return direct Word text and embedded images for OCR."""
    try:
        from docx import Document
    except ModuleNotFoundError:
        raise RuntimeError("Word .docx support needs python-docx. Install it with: pip install python-docx")

    data = uploaded_file.getvalue()
    doc = Document(io.BytesIO(data))

    text_parts = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if txt:
            text_parts.append(txt)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)

    images = []
    try:
        import zipfile
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            media_files = [n for n in zf.namelist() if n.startswith("word/media/")]
            for idx, name in enumerate(media_files[:int(max_images)], start=1):
                try:
                    img = Image.open(io.BytesIO(zf.read(name))).convert("RGB")
                    images.append((f"DOCX image {idx}", img))
                except Exception:
                    pass
    except Exception:
        pass

    return "\n".join(text_parts).strip(), images


def _extract_pptx_text_and_images(uploaded_file, max_images=25):
    """Return direct PowerPoint text and embedded images for OCR."""
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        raise RuntimeError("PowerPoint .pptx support needs python-pptx. Install it with: pip install python-pptx")

    data = uploaded_file.getvalue()
    prs = Presentation(io.BytesIO(data))

    text_parts = []
    images = []
    image_count = 0

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and str(shape.text).strip():
                    slide_text_parts.append(str(shape.text).strip())
            except Exception:
                pass

            if image_count < int(max_images):
                try:
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob)).convert("RGB")
                    image_count += 1
                    images.append((f"PPTX slide {slide_index} image {image_count}", img))
                except Exception:
                    pass

        if slide_text_parts:
            text_parts.append(f"--- PPTX Slide {slide_index} ---\n" + "\n".join(slide_text_parts))

    return "\n\n".join(text_parts).strip(), images


def _normalize_ocr_confidence(confidence):
    try:
        confidence = float(confidence)
        if confidence > 1.0:
            confidence = confidence / 100.0
        return max(0.0, min(1.0, confidence))
    except Exception:
        return 0.0

def page_ocr_predict(user):
    st.markdown("## 🔤 Text/Script OCR Prediction")
    st.caption("Upload an image, PDF, Word file, or PowerPoint file. Word/PPT typed text is extracted directly; images/pages are sent to OCR.")

    from backend.services.ocr_prediction_service import (
        get_available_ocr_models,
        load_ocr_model_from_path,
        predict_full_image_ocr,
    )

    uploaded = st.file_uploader(
        "Upload image / PDF / Word / PowerPoint",
        type=["png", "jpg", "jpeg", "bmp", "webp", "pdf", "docx", "pptx"],
        key="ocr_upload",
        help="Supported: images, PDF, Word .docx, PowerPoint .pptx",
    )

    if uploaded is None:
        st.info("Upload an image, PDF, Word (.docx), or PowerPoint (.pptx) file to start OCR/text extraction.")
        return

    file_ext = Path(uploaded.name).suffix.lower()
    image_exts = {".png", ".jpg", ".jpeg", ".bmp", ".webp"}

    direct_text = ""
    ocr_items = []   # list of (source_label, PIL image)

    try:
        if file_ext in image_exts:
            img = ImageOps.exif_transpose(Image.open(uploaded)).convert("RGB")
            ocr_items = [("Uploaded image", img)]
            st.image(img, caption="Uploaded image", use_container_width=True)

        elif file_ext == ".pdf":
            st.markdown('<div class="section-title">PDF Options</div>', unsafe_allow_html=True)
            max_pages = st.slider(
                "Maximum PDF pages to process",
                min_value=1,
                max_value=25,
                value=5,
                step=1,
                key="ocr_pdf_max_pages",
                help="Large PDFs can be slow. Increase only if needed.",
            )
            direct_text, ocr_items = _extract_pdf_text_and_images(uploaded, max_pages=max_pages, zoom=2.0)
            st.success(f"PDF loaded. Pages prepared for OCR: {len(ocr_items)}")
            if ocr_items:
                st.image(ocr_items[0][1], caption="PDF first page preview", use_container_width=True)

        elif file_ext == ".docx":
            st.markdown('<div class="section-title">Word Options</div>', unsafe_allow_html=True)
            max_images = st.slider(
                "Maximum embedded Word images to OCR",
                min_value=0,
                max_value=50,
                value=10,
                step=1,
                key="ocr_docx_max_images",
            )
            direct_text, ocr_items = _extract_docx_text_and_images(uploaded, max_images=max_images)
            st.success(f"Word file loaded. Embedded images prepared for OCR: {len(ocr_items)}")
            if direct_text:
                with st.expander("Preview extracted Word text", expanded=True):
                    st.text_area("Word text preview", direct_text[:5000], height=180, label_visibility="collapsed")
            if ocr_items:
                st.image(ocr_items[0][1], caption="First embedded Word image", use_container_width=True)

        elif file_ext == ".pptx":
            st.markdown('<div class="section-title">PowerPoint Options</div>', unsafe_allow_html=True)
            max_images = st.slider(
                "Maximum embedded PowerPoint images to OCR",
                min_value=0,
                max_value=75,
                value=15,
                step=1,
                key="ocr_pptx_max_images",
            )
            direct_text, ocr_items = _extract_pptx_text_and_images(uploaded, max_images=max_images)
            st.success(f"PowerPoint file loaded. Embedded images prepared for OCR: {len(ocr_items)}")
            if direct_text:
                with st.expander("Preview extracted PowerPoint text", expanded=True):
                    st.text_area("PowerPoint text preview", direct_text[:5000], height=180, label_visibility="collapsed")
            if ocr_items:
                st.image(ocr_items[0][1], caption="First embedded PowerPoint image", use_container_width=True)

        else:
            st.error("Unsupported file type. Please upload image, PDF, .docx, or .pptx.")
            return

    except Exception as e:
        st.error(f"File processing failed: {e}")
        st.info("For PDF/Word/PPT support install: python -m pip install pymupdf python-docx python-pptx")
        return

    # OCR model selector
    st.markdown('<div class="section-title">OCR Model</div>', unsafe_allow_html=True)
    ocr_model_map = get_available_ocr_models()
    ocr_labels = list(ocr_model_map.keys())

    if ocr_labels:
        selected_ocr_label = st.selectbox(
            "OCR Model",
            ocr_labels,
            index=0,
            key="ocr_page_model_select",
            label_visibility="collapsed",
        )
        ocr_model_path = str(ocr_model_map[selected_ocr_label])
    else:
        selected_ocr_label = None
        ocr_model_path = None
        if ocr_items:
            st.warning("No OCR model found in models/. Direct Word/PPT text can still be extracted, but image/page OCR needs a model or Printed Paragraph OCR engine.")

    st.markdown('<div class="section-title">OCR Reading Mode</div>', unsafe_allow_html=True)
    ocr_mode_label = st.selectbox(
        "OCR Reading Mode",
        [
    
            "Printed Paragraph / Screenshot",
            "Line / Words (recommended)",
            "Auto",
            "Character Grid - ",
            "Whole Image - Single Character",
        ],
        index=2,  # default: short word OCR model
        key="ocr_page_mode",
        label_visibility="collapsed",
        help="Choose how the app should read your image. Use Auto Detect if you are not sure.",
    )

    ocr_mode_map = {
        "Printed Text / Screenshot": "printed_paragraph",
        "Words or Line": "line",
        "Auto Detect": "auto",
        "Separate Characters": "character_grid",
        "Full Image": "whole",
    }

    force_upper = st.checkbox(
        "Convert final result to UPPERCASE",
        value=False,
        key="ocr_page_uppercase",
    )

    run_disabled = not direct_text and not ocr_items
    if st.button("Predict / Extract Text", type="primary", disabled=run_disabled):
        try:
            ocr_mode = ocr_mode_map.get(ocr_mode_label, "line")
            final_parts = []
            detail_rows = []

            if direct_text:
                extracted_text = direct_text.upper() if force_upper else direct_text
                final_parts.append("===== DIRECT DOCUMENT TEXT =====\n" + extracted_text)
                detail_rows.append({
                    "Source": "Document text",
                    "Type": "Direct extraction",
                    "Confidence (%)": "100.00",
                    "Text Preview": extracted_text[:120].replace("\n", " "),
                })

            loaded_ocr_model = None
            if ocr_items:
                external_ocr_modes = {"printed_paragraph"}
                if ocr_mode not in external_ocr_modes and not ocr_model_path:
                    st.error("Image/page OCR needs an OCR model in models/. Add your OCR .keras model, or choose Printed Paragraph / Handwritten Paragraph mode.")
                    return

                loaded_ocr_model = None if ocr_mode in external_ocr_modes else load_ocr_model_from_path(ocr_model_path)

                progress = st.progress(0, text="Running OCR...")
                total = len(ocr_items)

                for idx, (source_label, page_img) in enumerate(ocr_items, start=1):
                    progress.progress((idx - 1) / max(total, 1), text=f"OCR processing {source_label} ({idx}/{total})...")

                    result = predict_full_image_ocr(
                        page_img,
                        model=loaded_ocr_model,
                        mode=ocr_mode,
                    )

                    text = str(result.get("text", "") or "")
                    if force_upper:
                        text = text.upper()

                    confidence = _normalize_ocr_confidence(result.get("confidence", 0.0))
                    final_parts.append(f"===== OCR: {source_label} =====\n" + (text if text else "(no text detected)"))
                    detail_rows.append({
                        "Source": source_label,
                        "Type": f"OCR - {result.get('segmentation_mode', ocr_mode)}",
                        "Confidence (%)": f"{confidence * 100:.2f}",
                        "Text Preview": text[:120].replace("\n", " "),
                    })

                progress.progress(1.0, text="OCR complete.")

            final_text = "\n\n".join(final_parts).strip()
            if not final_text:
                final_text = "(no text detected)"

            # ── Save to universal history ──────────────────────
            # Calculate overall confidence roughly for history tracking
            avg_conf = 0.0
            ocr_rows = [r for r in detail_rows if r["Type"].startswith("OCR")]
            if ocr_rows:
                avg_conf = sum(float(r["Confidence (%)"]) for r in ocr_rows) / len(ocr_rows) / 100.0
            elif detail_rows:
                avg_conf = 1.0  # Direct extraction is 100%

            _render_result_and_accuracy(
                final_text,
                avg_conf * 100,
                result_label="Detected / Extracted Text",
            )

            _save_universal_prediction(
                category="script",
                user=user,
                input_source="document",
                source_filename=uploaded.name if uploaded else "",
                model_used="direct_extraction" if not ocr_items else (ocr_model_path or "external"),
                ocr_mode=ocr_mode if ocr_items else "text_extraction",
                ocr_text=final_text,
                ocr_confidence=avg_conf,
                line_count=len(final_text.splitlines()),
                extra_data={"detail_rows": detail_rows},
            )


            st.download_button(
                "⬇️ Download TXT",
                final_text,
                file_name=f"ocr_result_{_dt.datetime.now():%Y%m%d_%H%M%S}.txt",
                mime="text/plain",
                use_container_width=True,
            )


        except Exception as e:
            st.error(f"OCR prediction failed: {e}")
            if "Printed Paragraph mode needs EasyOCR or Tesseract" in str(e):
                st.info(
                    "Select **Line / Words (recommended)**, or install external OCR: "
                    "python -m pip install easyocr  OR install Tesseract and python -m pip install pytesseract"
                )

def page_client_dashboard(user):
    st.markdown(f"## 📊 My Dashboard — *{user['full_name']}*")
    try:
        SL = load_db()
        from backend.database.models import PredictionLog
        from sqlalchemy import func as sf
        db = SL()
        uid    = user["id"]
        bq     = db.query(PredictionLog).filter(PredictionLog.user_id == uid)
        total  = bq.count()
        avg_c  = db.query(sf.avg(PredictionLog.confidence)).filter(PredictionLog.user_id==uid).scalar() or 0
        avg_ms = db.query(sf.avg(PredictionLog.processing_time_ms)).filter(PredictionLog.user_id==uid).scalar() or 0
        mr     = bq.with_entities(PredictionLog.predicted_digit, sf.count(PredictionLog.id).label("c"))\
                    .group_by(PredictionLog.predicted_digit).order_by(sf.count(PredictionLog.id).desc()).first()
        most   = mr[0] if mr else "–"
        rows   = bq.order_by(PredictionLog.created_at.desc()).all()
        db.close()
    except Exception as e:
        st.error(f"Dashboard error: {e}"); return

    k1,k2,k3,k4 = st.columns(4)
    k1.metric("🧠 My Predictions", f"{total:,}")
    k2.metric("🎯 Avg Confidence",  f"{avg_c:.1f}%")
    k3.metric("⚡ Avg Response",     f"{avg_ms:.1f} ms")
    k4.metric("⭐ Most Predicted",   str(most))

    df = pd.DataFrame([{
        "predicted_digit": r.predicted_digit, "confidence": r.confidence,
        "model_used": r.model_used, "input_type": r.input_type,
        "processing_time_ms": r.processing_time_ms,
        "is_correct": r.is_correct, "created_at": r.created_at,
    } for r in rows])

    st.divider()
    if df.empty:
        st.info("No predictions yet — go to **Predict** to get started!"); return

    ca, cb = st.columns(2)
    with ca:
        st.markdown('<div class="section-title">Digit Distribution</div>', unsafe_allow_html=True)
        d = df["predicted_digit"].value_counts().sort_index().reset_index(); d.columns=["Digit","Count"]
        fig=px.bar(d,x="Digit",y="Count",color="Count",color_continuous_scale="Viridis",text="Count")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=260,margin=dict(t=10,b=0,l=0,r=0),
                          xaxis=dict(tickvals=list(range(10)),gridcolor="#2a2a3e"),
                          yaxis=dict(gridcolor="#2a2a3e"),coloraxis_showscale=False)
        st.plotly_chart(fig, use_container_width=True)

    with cb:
        st.markdown('<div class="section-title">Confidence Distribution</div>', unsafe_allow_html=True)
        fig=px.histogram(df,x="confidence",nbins=20,color_discrete_sequence=["#4a9eff"])
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=260,margin=dict(t=10,b=0,l=0,r=0),
                          xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig, use_container_width=True)

    cc, cd = st.columns(2)
    with cc:
        st.markdown('<div class="section-title">Model Usage</div>', unsafe_allow_html=True)
        mc=df["model_used"].value_counts().reset_index(); mc.columns=["Model","Count"]
        fig=px.bar(mc,x="Model",y="Count",color="Model",text="Count")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=260,margin=dict(t=10,b=0,l=0,r=0),
                          showlegend=False,xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig, use_container_width=True)

    with cd:
        st.markdown('<div class="section-title">Input Type Split</div>', unsafe_allow_html=True)
        it=df["input_type"].value_counts().reset_index(); it.columns=["Type","Count"]
        fig=px.pie(it,names="Type",values="Count",color_discrete_sequence=["#4a9eff","#a78bfa","#34d399"])
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",font_color="#c0c0e0",
                          height=260,margin=dict(t=10,b=0,l=40,r=40))
        st.plotly_chart(fig, use_container_width=True)

    if "created_at" in df.columns and df["created_at"].notna().any():
        st.markdown('<div class="section-title">Predictions Over Time</div>', unsafe_allow_html=True)
        df["date"] = pd.to_datetime(df["created_at"], utc=True).dt.date
        ts=df.groupby("date").size().reset_index(name="Count")
        fig=px.line(ts,x="date",y="Count",markers=True,color_discrete_sequence=["#4a9eff"])
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=240,margin=dict(t=10,b=0,l=0,r=0),
                          xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig, use_container_width=True)

    st.divider()
    st.markdown('<div class="section-title">📥 Export My Data</div>', unsafe_allow_html=True)
    stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    dex = df.copy(); dex["created_at"] = dex["created_at"].apply(_to_ist)
    ec1,ec2 = st.columns(2)
    with ec1:
        st.download_button("⬇️ Download CSV", dex.to_csv(index=False),
                           f"my_history_{stamp}.csv","text/csv",use_container_width=True)
    with ec2:
        st.download_button("📄 Download PDF", _personal_pdf(dex, user["full_name"]),
                           f"my_report_{stamp}.pdf","application/pdf",use_container_width=True)


def _personal_pdf(df, name):
    from io import BytesIO
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    buf=BytesIO()
    doc=SimpleDocTemplate(buf,pagesize=landscape(A4),leftMargin=1*cm,rightMargin=1*cm,
                           topMargin=1.5*cm,bottomMargin=1.5*cm)
    styl=getSampleStyleSheet()
    cols=[c for c in ["predicted_digit","confidence","model_used","input_type",
                       "is_correct","processing_time_ms","created_at"] if c in df.columns]
    sub=df[cols].copy(); sub.columns=[c.replace("_"," ").title() for c in cols]
    data=[list(sub.columns)]+sub.astype(str).values.tolist()
    tbl=Table(data,repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#1E3A8A")),
        ("TEXTCOLOR",(0,0),(-1,0),colors.white),
        ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),
        ("FONTSIZE",(0,0),(-1,-1),7.5),
        ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#EFF6FF")]),
        ("GRID",(0,0),(-1,-1),.4,colors.HexColor("#CBD5E1")),
        ("ALIGN",(0,0),(-1,-1),"CENTER"),
        ("VALIGN",(0,0),(-1,-1),"MIDDLE"),
        ("TOPPADDING",(0,0),(-1,-1),4),("BOTTOMPADDING",(0,0),(-1,-1),4),
    ]))
    doc.build([
        Paragraph(f"<b>Personal History — {name}</b>",styl["Title"]),
        Paragraph(f"Generated: {_dt.datetime.now():%Y-%m-%d %H:%M} | Records: {len(df)}",styl["Normal"]),
        Spacer(1,.4*cm), tbl,
    ])
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: CLIENT HISTORY
# ══════════════════════════════════════════════════════════════════════════════
def page_client_history(user):
    st.markdown("## 🕐 My Prediction History")
    
    cat = st.radio("Category", ["Digit", "Multi-Digit", "Alphabet", "Script"], horizontal=True)
    f1, f2, f3 = st.columns([2, 2, 1])
    with f1: mf = st.selectbox("Model", ["All Models"] + AVAILABLE_MODELS)
    with f2:
        if cat == "Digit":
            df_ = st.selectbox("Digit", ["All Digits"] + [str(i) for i in range(10)])
        else:
            df_ = "All Digits"
    with f3: pp = st.selectbox("Rows/page", [10, 20, 50], index=1)

    try:
        SL = load_db()
        from backend.database.models import PredictionLog, UniversalPredictionLog
        db = SL()
        
        rows = []
        tot = 0
        tp = 1
        
        if "hist_pg" not in st.session_state: st.session_state["hist_pg"] = 1
        pg = st.session_state["hist_pg"]
        
        if cat == "Digit":
            # Combine legacy PredictionLog and new UniversalPredictionLog
            q1 = db.query(PredictionLog).filter(PredictionLog.user_id == user["id"])
            if mf != "All Models": q1 = q1.filter(PredictionLog.model_used == mf)
            if df_ != "All Digits": q1 = q1.filter(PredictionLog.predicted_digit == int(df_))
            recs1 = q1.order_by(PredictionLog.id.desc()).all()
            
            q2 = db.query(UniversalPredictionLog).filter(
                UniversalPredictionLog.user_id == user["id"],
                UniversalPredictionLog.prediction_category == "digit"
            )
            if mf != "All Models": q2 = q2.filter(UniversalPredictionLog.model_used == mf)
            if df_ != "All Digits": q2 = q2.filter(UniversalPredictionLog.predicted_digit == int(df_))
            recs2 = q2.order_by(UniversalPredictionLog.id.desc()).all()
            
            # Merge and sort in memory since total count for a single user is usually small
            all_recs = sorted(recs1 + recs2, key=lambda x: x.created_at, reverse=True)
            tot = len(all_recs)
            tp = max(1, (tot + pp - 1) // pp)
            pg = max(1, min(pg, tp))
            
            page_recs = all_recs[(pg - 1) * pp : pg * pp]
            
            for r in page_recs:
                t3 = " | ".join(f"{t['digit']}:{t['confidence']:.1f}%" for t in (r.top3_predictions or [])[:3])
                rows.append({
                    "Date": _to_ist(r.created_at),
                    "Digit": r.predicted_digit,
                    "Conf %": round(r.confidence, 1) if r.confidence else 0,
                    "Top-3": t3,
                    "Model": r.model_used,
                    "Type": getattr(r, "input_type", getattr(r, "input_source", "unknown")),
                })
                
        else:
            cat_map = {"Multi-Digit": "multidigit", "Alphabet": "alphabet", "Script": "script"}
            db_cat = cat_map[cat]
            
            q = db.query(UniversalPredictionLog).filter(
                UniversalPredictionLog.user_id == user["id"],
                UniversalPredictionLog.prediction_category == db_cat
            )
            if mf != "All Models": q = q.filter(UniversalPredictionLog.model_used == mf)
            
            tot = q.count()
            tp = max(1, (tot + pp - 1) // pp)
            pg = max(1, min(pg, tp))
            
            recs = q.order_by(UniversalPredictionLog.id.desc()).offset((pg - 1) * pp).limit(pp).all()
            
            for r in recs:
                if db_cat == "multidigit":
                    rows.append({
                        "Date": _to_ist(r.created_at),
                        "Result": r.multidigit_result,
                        "Count": r.digit_count,
                        "Model": r.model_used,
                        "Source": r.input_source,
                    })
                elif db_cat == "alphabet":
                    rows.append({
                        "Date": _to_ist(r.created_at),
                        "Extracted Text": r.ocr_text,
                        "Conf %": round((r.ocr_confidence or 0) * 100, 1),
                        "Mode": r.ocr_mode,
                        "Source": r.input_source,
                    })
                elif db_cat == "script":
                    rows.append({
                        "Date": _to_ist(r.created_at),
                        "Document": r.source_filename or "Image/Text",
                        "Extracted Text": (r.ocr_text[:50] + "...") if r.ocr_text and len(r.ocr_text) > 50 else r.ocr_text,
                        "Conf %": round((r.ocr_confidence or 0) * 100, 1),
                        "Lines": r.line_count,
                        "Mode": r.ocr_mode,
                    })

        db.close()
        
        if rows:
            dfs = pd.DataFrame(rows)
            st.dataframe(dfs, use_container_width=True, hide_index=True, height=420)
            
            pc = st.columns([1, 3, 1])
            with pc[0]:
                if st.button("◀ Prev", disabled=(pg <= 1), key="cp_prev"): 
                    st.session_state["hist_pg"] = pg - 1
                    st.rerun()
            with pc[1]:
                st.markdown(f"<p style='text-align:center;color:#7070a0;margin-top:.4rem'>Page {pg}/{tp} | {tot} records</p>", unsafe_allow_html=True)
            with pc[2]:
                if st.button("Next ▶", disabled=(pg >= tp), key="cp_next"): 
                    st.session_state["hist_pg"] = pg + 1
                    st.rerun()
                    
            stamp = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
            st.download_button("⬇️ CSV", dfs.to_csv(index=False), f"{cat.lower()}_hist_{stamp}.csv", "text/csv")
        else:
            st.info(f"No {cat} records found.")
            
    except Exception as e:
        st.error(f"History error: {e}")




# ══════════════════════════════════════════════════════════════════════════════
# PAGE: ADMIN DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
def page_admin_dashboard():
    st.markdown("## 📊 Admin Analytics Dashboard")
    try:
        SL=load_db()
        from backend.database.models import PredictionLog,User
        from sqlalchemy import func as sf
        db=SL()
        total=db.query(sf.count(PredictionLog.id)).scalar() or 0
        avg_c=db.query(sf.avg(PredictionLog.confidence)).scalar() or 0
        avg_ms=db.query(sf.avg(PredictionLog.processing_time_ms)).scalar() or 0
        mr=(db.query(PredictionLog.predicted_digit,sf.count(PredictionLog.id).label("c"))
              .group_by(PredictionLog.predicted_digit).order_by(sf.count(PredictionLog.id).desc()).first())
        tu=db.query(sf.count(User.id)).scalar() or 0
        md=mr[0] if mr else "–"
        db.close()
    except Exception as e: st.error(f"Error: {e}"); return

    k1,k2,k3,k4,k5=st.columns(5)
    k1.metric("🧠 Total Predictions",f"{total:,}")
    k2.metric("🎯 Avg Confidence",f"{avg_c:.1f}%")
    k3.metric("⚡ Avg Response",f"{avg_ms:.1f} ms")
    k4.metric("⭐ Most Predicted",str(md))
    k5.metric("👥 Total Users",f"{tu:,}")
    st.divider()

    try:
        SL=load_db()
        from backend.database.models import PredictionLog
        db=SL()
        all_rows=db.query(PredictionLog).all(); db.close()
        df=pd.DataFrame([{"username":r.username or "unknown","predicted_digit":r.predicted_digit,
                           "confidence":r.confidence,"model_used":r.model_used,
                           "input_type":r.input_type,"processing_time_ms":r.processing_time_ms,
                           "is_correct":r.is_correct,"created_at":r.created_at} for r in all_rows])
    except: df=pd.DataFrame()

    if df.empty: st.info("No data yet."); return

    ca,cb=st.columns(2)
    with ca:
        st.markdown('<div class="section-title">Global Digit Distribution</div>',unsafe_allow_html=True)
        d=df["predicted_digit"].value_counts().sort_index().reset_index(); d.columns=["Digit","Count"]
        fig=px.bar(d,x="Digit",y="Count",color="Count",color_continuous_scale="Viridis",text="Count")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=270,margin=dict(t=10,b=0,l=0,r=0),
                          xaxis=dict(tickvals=list(range(10)),gridcolor="#2a2a3e"),
                          yaxis=dict(gridcolor="#2a2a3e"),coloraxis_showscale=False)
        st.plotly_chart(fig,use_container_width=True)

    with cb:
        st.markdown('<div class="section-title">Confidence Histogram</div>',unsafe_allow_html=True)
        fig=px.histogram(df,x="confidence",nbins=20,color_discrete_sequence=["#4a9eff"])
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=270,margin=dict(t=10,b=0,l=0,r=0),
                          xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig,use_container_width=True)

    cc,cd=st.columns(2)
    with cc:
        st.markdown('<div class="section-title">Model Usage</div>',unsafe_allow_html=True)
        mc=df["model_used"].value_counts().reset_index(); mc.columns=["Model","Count"]
        fig=px.bar(mc,x="Model",y="Count",color="Model",text="Count")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=270,margin=dict(t=10,b=0,l=0,r=0),
                          showlegend=False,xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig,use_container_width=True)

    with cd:
        st.markdown('<div class="section-title">Per-User Predictions</div>',unsafe_allow_html=True)
        uc=df["username"].value_counts().reset_index(); uc.columns=["User","Count"]
        fig=px.bar(uc,x="User",y="Count",color="Count",color_continuous_scale="Plasma",text="Count")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=270,margin=dict(t=10,b=0,l=0,r=0),
                          coloraxis_showscale=False,xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e"))
        st.plotly_chart(fig,use_container_width=True)

    # Model accuracy
    st.divider()
    st.markdown('<div class="section-title">📈 Model Accuracy Comparison</div>',unsafe_allow_html=True)
    mrows=[]
    seen=set()
    for p in sorted(MODELS_DIR.glob("*_metrics.json")):
        if p.name.endswith("_full_metrics.json"): continue
        try:
            data=json.loads(p.read_text())
            mn=data.get("model",p.stem.replace("_metrics",""))
            ra=data.get("accuracy",data.get("test_accuracy",0)) or 0
            ap=round(ra*100,2) if ra<=1.0 else round(ra,2)
            if mn in seen:
                ex=next(r for r in mrows if r["Model"]==mn)
                if ap>ex["Accuracy"]: mrows.remove(ex)
                else: continue
            seen.add(mn); mrows.append({"Model":mn,"Accuracy":ap})
        except: pass
    if mrows:
        dfm=pd.DataFrame(mrows).sort_values("Accuracy",ascending=False).reset_index(drop=True)
        fig=px.bar(dfm,x="Model",y="Accuracy",color="Accuracy",color_continuous_scale="RdYlGn",
                   range_y=[max(80,dfm["Accuracy"].min()-2),100.5],
                   text=dfm["Accuracy"].apply(lambda v:f"{v:.2f}%"))
        fig.update_traces(textposition="outside")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=380,margin=dict(t=30,b=80,l=0,r=0),
                          xaxis=dict(gridcolor="#2a2a3e",tickangle=-35),
                          yaxis=dict(gridcolor="#2a2a3e"),coloraxis_showscale=False)
        st.plotly_chart(fig,use_container_width=True)

    # Export
    st.divider()
    st.markdown('<div class="section-title">📂 Export All Data</div>',unsafe_allow_html=True)
    try:
        from backend.services.report_service import generate_pdf_report, generate_csv, generate_excel
        recs=df.to_dict("records")
        tp=len(recs); ac=sum(r["confidence"] or 0 for r in recs)/max(tp,1)
        at=sum(r["processing_time_ms"] or 0 for r in recs)/max(tp,1)
        dc={};
        for r in recs: d=r["predicted_digit"]; dc[d]=dc.get(d,0)+1
        summary={"total":tp,"avg_confidence":round(ac,2),"avg_time_ms":round(at,2),
                 "canvas_count":sum(1 for r in recs if r.get("input_type")=="canvas"),
                 "upload_count":sum(1 for r in recs if r.get("input_type")=="upload"),
                 "most_common_digit":max(dc,key=dc.get) if dc else None}
        stamp=_dt.datetime.now().strftime("%Y%m%d_%H%M%S")
        rc1,rc2,rc3=st.columns(3)
        with rc1: st.download_button("📄 PDF Report",generate_pdf_report(recs,summary),
                                     f"admin_report_{stamp}.pdf","application/pdf",use_container_width=True)
        with rc2: st.download_button("📊 CSV",generate_csv(recs),
                                     f"all_pred_{stamp}.csv","text/csv",use_container_width=True)
        with rc3: st.download_button("📗 Excel",generate_excel(recs),
                                     f"all_pred_{stamp}.xlsx",
                                     "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                     use_container_width=True)
    except Exception as e: st.error(f"Export error: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: ADMIN ALL HISTORY
# ══════════════════════════════════════════════════════════════════════════════
# PAGE: ADMIN ALL HISTORY
# ══════════════════════════════════════════════════════════════════════════════
def page_admin_history():
    st.markdown("## 🕐 All Users — Prediction History")
    try:
        SL=load_db()
        from backend.database.models import User
        db=SL()
        unames=["All Users"]+[u.username for u in db.query(User).order_by(User.username).all()]
        db.close()
    except: unames=["All Users"]

    cat = st.radio("Category", ["Digit", "Multi-Digit", "Alphabet", "Script"], horizontal=True, key="ah_cat")
    
    f1,f2,f3,f4=st.columns([2,2,2,1])
    with f1: uf=st.selectbox("User",unames)
    with f2: mf=st.selectbox("Model",["All Models"]+AVAILABLE_MODELS,key="ah_mod")
    with f3: 
        if cat == "Digit":
            df_=st.selectbox("Digit",["All Digits"]+[str(i) for i in range(10)],key="ah_dig")
        else:
            df_="All Digits"
    with f4: pp=st.selectbox("Rows",[20,50,100],key="ah_pp")

    try:
        SL=load_db()
        from backend.database.models import PredictionLog, UniversalPredictionLog
        db=SL()
        
        rows = []
        tot = 0
        tp = 1
        
        if "ahp" not in st.session_state: st.session_state["ahp"]=1
        pg = st.session_state["ahp"]

        if cat == "Digit":
            q1=db.query(PredictionLog)
            if uf!="All Users":   q1=q1.filter(PredictionLog.username==uf)
            if mf!="All Models":  q1=q1.filter(PredictionLog.model_used==mf)
            if df_!="All Digits": q1=q1.filter(PredictionLog.predicted_digit==int(df_))
            recs1 = q1.order_by(PredictionLog.id.desc()).all()
            
            q2=db.query(UniversalPredictionLog).filter(UniversalPredictionLog.prediction_category == "digit")
            if uf!="All Users":   q2=q2.filter(UniversalPredictionLog.username==uf)
            if mf!="All Models":  q2=q2.filter(UniversalPredictionLog.model_used==mf)
            if df_!="All Digits": q2=q2.filter(UniversalPredictionLog.predicted_digit==int(df_))
            recs2 = q2.order_by(UniversalPredictionLog.id.desc()).all()
            
            all_recs = sorted(recs1 + recs2, key=lambda x: x.created_at, reverse=True)
            tot = len(all_recs)
            tp = max(1, (tot+pp-1)//pp)
            pg = max(1, min(pg, tp))
            page_recs = all_recs[(pg-1)*pp : pg*pp]
            
            for r in page_recs:
                t3=" | ".join(f"{t['digit']}:{t['confidence']:.1f}%" for t in (r.top3_predictions or [])[:3])
                rows.append({"User":r.username or "–","Date":_to_ist(r.created_at),
                             "Digit":r.predicted_digit,
                             "Conf %":round(r.confidence,1) if r.confidence else 0,
                             "Top-3":t3,"Model":r.model_used,
                             "Type":getattr(r, "input_type", getattr(r, "input_source", "unknown"))})
        else:
            cat_map = {"Multi-Digit": "multidigit", "Alphabet": "alphabet", "Script": "script"}
            db_cat = cat_map[cat]
            
            q = db.query(UniversalPredictionLog).filter(UniversalPredictionLog.prediction_category == db_cat)
            if uf!="All Users":   q=q.filter(UniversalPredictionLog.username==uf)
            if mf!="All Models":  q=q.filter(UniversalPredictionLog.model_used==mf)
            
            tot = q.count()
            tp = max(1, (tot+pp-1)//pp)
            pg = max(1, min(pg, tp))
            
            recs = q.order_by(UniversalPredictionLog.id.desc()).offset((pg-1)*pp).limit(pp).all()
            
            for r in recs:
                if db_cat == "multidigit":
                    rows.append({
                        "User": r.username or "–",
                        "Date": _to_ist(r.created_at),
                        "Result": r.multidigit_result,
                        "Count": r.digit_count,
                        "Model": r.model_used,
                        "Source": r.input_source,
                    })
                elif db_cat == "alphabet":
                    rows.append({
                        "User": r.username or "–",
                        "Date": _to_ist(r.created_at),
                        "Extracted Text": r.ocr_text,
                        "Conf %": round((r.ocr_confidence or 0) * 100, 1),
                        "Mode": r.ocr_mode,
                        "Source": r.input_source,
                    })
                elif db_cat == "script":
                    rows.append({
                        "User": r.username or "–",
                        "Date": _to_ist(r.created_at),
                        "Document": r.source_filename or "Image/Text",
                        "Extracted Text": (r.ocr_text[:50] + "...") if r.ocr_text and len(r.ocr_text) > 50 else r.ocr_text,
                        "Conf %": round((r.ocr_confidence or 0) * 100, 1),
                        "Lines": r.line_count,
                        "Mode": r.ocr_mode,
                    })

        db.close()
        
        if rows:
            dfs=pd.DataFrame(rows)
            st.dataframe(dfs,use_container_width=True,hide_index=True,height=480)
            pc=st.columns([1,3,1])
            with pc[0]:
                if st.button("◀ Prev",key="ahprev",disabled=(pg<=1)): st.session_state["ahp"]=pg-1; st.rerun()
            with pc[1]:
                st.markdown(f"<p style='text-align:center;color:#7070a0;margin-top:.4rem'>Page {pg}/{tp} | {tot} records</p>",unsafe_allow_html=True)
            with pc[2]:
                if st.button("Next ▶",key="ahnext",disabled=(pg>=tp)): st.session_state["ahp"]=pg+1; st.rerun()
            stamp=_dt.datetime.now().strftime("%Y%m%d_%H%M%S")
            st.download_button("⬇️ Download CSV",dfs.to_csv(index=False),f"all_{cat.lower()}_hist_{stamp}.csv","text/csv")
        else:
            st.info(f"No {cat} records found.")
    except Exception as e: st.error(f"History error: {e}")


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: USER MANAGEMENT (admin)
# ══════════════════════════════════════════════════════════════════════════════
def page_user_management():
    st.markdown("## 👥 User Management")
    try:
        SL=load_db()
        from backend.database.models import User, PredictionLog
        from sqlalchemy import func as sf
        db=SL()
        users=db.query(User).order_by(User.created_at.desc()).all()
        pcs=dict(db.query(PredictionLog.user_id,sf.count(PredictionLog.id)).group_by(PredictionLog.user_id).all())
        db.close()
        st.dataframe(pd.DataFrame([{
            "ID":u.id,"Username":u.username,"Full Name":u.full_name or "–",
            "Email":u.email or "–","Role":u.role,
            "Active":"✅" if u.is_active else "❌",
            "Predictions":pcs.get(u.id,0),
            "Created":_to_ist(u.created_at),"Last Login":_to_ist(u.last_login),
        } for u in users]), use_container_width=True, hide_index=True)
    except Exception as e: st.error(f"Error: {e}")

    st.divider()
    st.markdown('<div class="section-title">Create New User</div>',unsafe_allow_html=True)
    with st.form("cu_form"):
        c1,c2,c3=st.columns(3)
        with c1: nn=st.text_input("Full Name"); nu=st.text_input("Username")
        with c2: ne=st.text_input("Email"); nr=st.selectbox("Role",["client","admin"])
        with c3: np1=st.text_input("Password",type="password",help="Min 8 · 1 UPPER · 1 special"); np2=st.text_input("Confirm",type="password")
        if st.form_submit_button("➕ Create User",use_container_width=True):
            if not nn or not nu or not np1: st.error("Name, username, password required.")
            elif np1!=np2: st.error("Passwords don't match.")
            else:
                err=_validate_password(np1)
                if err: st.error(err)
                else:
                    try:
                        SL=load_db()
                        from backend.database.models import User
                        db=SL()
                        if db.query(User).filter(User.username==nu).first():
                            st.error("Username taken.")
                        else:
                            db.add(User(username=nu,password_hash=_hash_pw(np1),role=nr,full_name=nn,email=ne))
                            db.commit(); db.close(); st.success(f"User **{nu}** created!"); st.rerun()
                    except Exception as e: st.error(f"Error: {e}")

    st.divider()
    st.markdown('<div class="section-title">Delete User & Data</div>',unsafe_allow_html=True)
    with st.form("du_form"):
        st.warning("⚠️ This will permanently delete the user account, their history, and all associated data. This cannot be undone.")
        
        current_admin = st.session_state.get("auth_user", {}).get("username")
        try:
            SL=load_db()
            db=SL()
            from backend.database.models import User, PredictionLog, OTPRecord
            all_users = [u.username for u in db.query(User).order_by(User.username).all() if u.username != current_admin]
        except Exception:
            all_users = []
        finally:
            if 'db' in locals(): db.close()

        u_del = st.selectbox("Select User to Delete", [""] + all_users)
        
        if st.form_submit_button("🗑️ Permanently Delete User", use_container_width=True):
            if not u_del: 
                st.error("Please select a user to delete.")
            else:
                try:
                    SL=load_db()
                    db=SL()
                    target = db.query(User).filter(User.username == u_del).first()
                    if target:
                        db.query(OTPRecord).filter(OTPRecord.user_id == target.id).delete()
                        db.query(PredictionLog).filter(PredictionLog.user_id == target.id).delete()
                        db.delete(target)
                        db.commit()
                        st.success(f"User **{u_del}** and all associated data were successfully deleted.")
                        st.rerun()
                    else:
                        st.error("User not found.")
                except Exception as e:
                    db.rollback()
                    st.error(f"Error deleting user: {e}")
                finally:
                    if 'db' in locals(): db.close()



# ══════════════════════════════════════════════════════════════════════════════
# PAGE: ACCURACY (admin)
# ══════════════════════════════════════════════════════════════════════════════
def page_accuracy_admin():
    st.markdown("## 📈 Model Accuracy & Performance")
    mrows=[]; seen=set()
    for p in sorted(MODELS_DIR.glob("*_metrics.json")):
        if p.name.endswith("_full_metrics.json"): continue
        try:
            data=json.loads(p.read_text())
            mn=data.get("model",p.stem.replace("_metrics",""))
            ra=data.get("accuracy",data.get("test_accuracy",0)) or 0
            ap=round(ra*100,2) if ra<=1.0 else round(ra,2)
            rf=data.get("macro_f1",0) or 0
            fp=round(rf*100,2) if rf<=1.0 else round(rf,2)
            ro=data.get("auc_roc",0) or 0
            op=round(ro*100,2) if ro<=1.0 else round(ro,2)
            if mn in seen:
                ex=next(r for r in mrows if r["Model"]==mn)
                if ap>ex["Accuracy"]: mrows.remove(ex)
                else: continue
            seen.add(mn); mrows.append({"Model":mn,"Accuracy":ap,"Macro F1":fp,"AUC-ROC":op})
        except: pass
    if mrows:
        dfm=pd.DataFrame(mrows).sort_values("Accuracy",ascending=False).reset_index(drop=True)
        fig=px.bar(dfm,x="Model",y="Accuracy",color="Accuracy",color_continuous_scale="RdYlGn",
                   range_y=[max(80,dfm["Accuracy"].min()-2),100.5],
                   text=dfm["Accuracy"].apply(lambda v:f"{v:.2f}%"))
        fig.update_traces(textposition="outside")
        fig.update_layout(paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                          font_color="#c0c0e0",height=400,margin=dict(t=30,b=80,l=0,r=0),
                          xaxis=dict(gridcolor="#2a2a3e",tickangle=-35),
                          yaxis=dict(gridcolor="#2a2a3e"),coloraxis_showscale=False)
        st.plotly_chart(fig,use_container_width=True)

        st.markdown('<div class="section-title">Multi-Metric Comparison</div>',unsafe_allow_html=True)
        fig2=go.Figure()
        for _,row in dfm.iterrows():
            fig2.add_trace(go.Bar(name=row["Model"],x=["Accuracy","Macro F1","AUC-ROC"],
                                  y=[row["Accuracy"],row["Macro F1"],row["AUC-ROC"]]))
        fig2.update_layout(barmode="group",paper_bgcolor="rgba(0,0,0,0)",plot_bgcolor="rgba(0,0,0,0)",
                           font_color="#c0c0e0",height=380,margin=dict(t=10,b=0,l=0,r=0),
                           xaxis=dict(gridcolor="#2a2a3e"),yaxis=dict(gridcolor="#2a2a3e",range=[0,105]))
        st.plotly_chart(fig2,use_container_width=True)

        try: st.dataframe(dfm.style.background_gradient(subset=["Accuracy"],cmap="RdYlGn"),use_container_width=True,hide_index=True)
        except: st.dataframe(dfm,use_container_width=True,hide_index=True)
        st.download_button("⬇️ Metrics CSV",dfm.to_csv(index=False),"model_metrics.csv","text/csv")
    else:
        st.info("No model metrics files found in models/.")



# ══════════════════════════════════════════════════════════════════════════════
# PAGE: CANVAS PREDICTION - MERGED CANVAS FEATURES
# ══════════════════════════════════════════════════════════════════════════════
def _canvas_has_ink(canvas_array, background="white"):
    """Return True only when the user actually drew something on the canvas."""
    try:
        arr = np.asarray(canvas_array).astype(np.uint8)
        rgb = arr[:, :, :3] if arr.ndim == 3 else arr
        if background == "white":
            return bool(np.any(rgb < 245))
        return bool(np.any(rgb > 10))
    except Exception:
        return False


def _uppercase_ocr_result_inplace(ocr_result):
    """Convert OCR result text, line text, and character text to uppercase."""
    ocr_result["text"] = str(ocr_result.get("text", "")).upper()
    for line in ocr_result.get("lines", []):
        if "text" in line:
            line["text"] = str(line.get("text", "")).upper()
        for ch in line.get("char_details", []):
            if "char" in ch:
                ch["char"] = str(ch.get("char", "")).upper()
    return ocr_result



def _render_ocr_result_block(ocr_result):
    """
    Compact OCR result renderer.
    Only shows predicted text and accuracy/confidence.
    """
    predicted_text = str(ocr_result.get("text", "") or "")
    confidence = ocr_result.get("confidence", 0.0)
    _render_result_and_accuracy(
        predicted_text if predicted_text else "(no text detected)",
        confidence,
        result_label="Predicted Text",
    )



def _render_universal_digit_result(result, show_debug=False):
    """
    Compact multi-digit result renderer.
    Only shows final number/result and average accuracy.
    """
    final_number_output = result.get("number_prediction") or result.get("prediction", "")

    conf_values = []
    for d in result.get("digits", []) or []:
        c = d.get("confidence", None)
        pct = _confidence_to_percent(c)
        if pct is not None:
            conf_values.append(pct)

    avg_conf = sum(conf_values) / len(conf_values) if conf_values else None

    _render_result_and_accuracy(
        final_number_output if final_number_output else "(no digit detected)",
        avg_conf,
        result_label="Predicted Number",
    )


def _render_clear_steps(mode="canvas"):
    # Compact UI: no large step cards.
    return


def _render_choice_help(selected, kind="canvas"):
    # Compact UI: no explanatory cards.
    return


def _clear_intro(title, body):
    # Compact UI: no large workflow panel.
    return


def _clear_section(title, body):
    st.caption(body)




def _compact_choice(label, options, key, default=None):
    """Modern selector with fallback for older Streamlit versions."""
    if hasattr(st, "segmented_control"):
        return st.segmented_control(label, options, default=default or options[0], key=key)
    return st.radio(label, options, index=options.index(default or options[0]), horizontal=True, key=key)



def _validate_canvas_digit_like(image_arr, mode="single_digit"):
    """
    Reject invalid canvas drawings before sending them to the digit CNN.

    This prevents forced predictions for:
    - single dots / tiny marks
    - very small scribbles
    - straight horizontal or vertical lines
    - thin zig-zag lines that do not look like digits
    """
    import cv2
    import numpy as np

    if image_arr is None:
        return False, "Please draw something first."

    arr = np.array(image_arr)

    # Remove alpha channel if canvas returns RGBA
    if arr.ndim == 3 and arr.shape[2] == 4:
        arr = arr[:, :, :3]

    # Convert to grayscale
    if arr.ndim == 3:
        gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    else:
        gray = arr.copy()

    # Canvas uses black drawing on white background.
    # Foreground = dark pixels.
    _, mask = cv2.threshold(gray, 230, 255, cv2.THRESH_BINARY_INV)

    # Remove tiny noise
    kernel = np.ones((3, 3), np.uint8)
    mask = cv2.morphologyEx(mask, cv2.MORPH_OPEN, kernel)

    ink_pixels = cv2.countNonZero(mask)
    total_pixels = mask.shape[0] * mask.shape[1]
    ink_ratio = ink_pixels / max(total_pixels, 1)

    # Too little ink means dot/noise/incomplete drawing
    if ink_pixels < 120:
        return False, "Drawing is too small. Please draw a complete digit."

    if ink_ratio < 0.003:
        return False, "Very little ink detected. Please draw a clear digit."

    contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    if not contours:
        return False, "No valid digit shape detected."

    largest = max(contours, key=cv2.contourArea)
    area = cv2.contourArea(largest)
    x, y, w, h = cv2.boundingRect(largest)
    bbox_area = w * h
    aspect_ratio = w / max(h, 1)

    # Reject tiny bounding boxes
    if w < 18 or h < 28:
        return False, "Drawing is too small. Please draw a larger digit."

    if mode == "single_digit":
        if aspect_ratio > 2.2:
            return False, "This looks like a horizontal line or scribble, not a single digit."

        if aspect_ratio < 0.25:
            return False, "This looks like a vertical line, not a complete digit."

    if mode == "multi_digit":
        if h < 35:
            return False, "Drawing height is too small. Please draw clearer digits."

    fill_ratio = area / max(bbox_area, 1)
    if fill_ratio < 0.015:
        return False, "This looks like a thin line, not a digit."

    perimeter = cv2.arcLength(largest, True)
    if perimeter < 60:
        return False, "Drawing is too simple. Please draw a complete digit."

    rect = cv2.minAreaRect(largest)
    rw, rh = rect[1]
    if rw > 0 and rh > 0:
        line_ratio = max(rw, rh) / max(min(rw, rh), 1)
        if line_ratio > 12:
            return False, "This looks like a straight line, not a digit."

    return True, "Valid digit-like drawing."


def page_canvas_prediction(user):
    st.markdown("""
    <div class="page-hero">
      <span class="mode-chip">🖊️ Canvas</span>
      <h2>Canvas Prediction</h2>
      <p>Draw a digit, number, alphabet, or word and run prediction.</p>
    </div>
    """, unsafe_allow_html=True)

    prediction_type = _compact_choice(
        "Choose prediction type",
        ["Single Digit", "Multi-Digit", "OCR Text / Alphabet"],
        default="Single Digit",
        key="canvas_prediction_type",
    )

    st.divider()

    # ──────────────────────────────────────────────────────────────────────
    # 1. SINGLE DIGIT CANVAS
    # ──────────────────────────────────────────────────────────────────────
    if prediction_type == "Single Digit":
        st.markdown("### 🔢 Single Digit Prediction from Canvas")
        _clear_section("Use this for one digit only", "Example: draw only `7`, not `78`. Keep the digit centered inside the canvas for better accuracy.")
        col_left, col_right = st.columns([1, 1], gap="large")

        with col_left:
            st.markdown('<div class="section-title">Draw one digit 0–9</div>', unsafe_allow_html=True)
            brush = st.slider("Brush size", 8, 30, 16, key="single_digit_canvas_brush")
            canvas_result = st_canvas(
                fill_color="rgba(255,255,255,0)",
                stroke_width=brush,
                stroke_color="#000000",
                background_color="#FFFFFF",
                height=280,
                width=280,
                drawing_mode="freedraw",
                key="single_digit_canvas",
            )

            input_image = None
            if canvas_result.image_data is not None:
                arr = canvas_result.image_data.astype(np.uint8)
                if _canvas_has_ink(arr, background="white"):
                    input_image = arr

            st.markdown('<div class="section-title">Model & Options</div>', unsafe_allow_html=True)
            sel_model = st.selectbox(
                "Model",
                AVAILABLE_MODELS,
                format_func=lambda m: MODEL_LABELS.get(m, m),
                index=min(1, len(AVAILABLE_MODELS) - 1),
                key="single_digit_canvas_model",
            )
            xai_on = st.checkbox("Show XAI heatmaps", value=True, key="single_digit_canvas_xai")
            pred_btn = st.button(
                "🔮 Predict Single Digit",
                type="primary",
                use_container_width=True,
                disabled=(input_image is None),
                key="single_digit_canvas_btn",
            )

        with col_right:
            if pred_btn and input_image is not None:
                valid, msg = _validate_canvas_digit_like(input_image, mode="single_digit")
                if not valid:
                    st.warning(f"⚠️ {msg}")
                    return

                with st.spinner("Running single digit prediction..."):
                    try:
                        pfn, gcfn, limefn, fgsmfn, reg = load_prediction_service()
                        result = pfn(input_image, sel_model)
                        st.session_state.update({
                            "last_result": result,
                            "last_canvas_single_result": result,
                            "last_image": input_image,
                            "last_model": sel_model,
                            "last_input_type": "canvas_single_digit",
                            "xai_on": xai_on,
                        })
                        if xai_on:
                            st.session_state["gradcam_path"] = gcfn(input_image, sel_model)
                            st.session_state["lime_path"] = limefn(input_image, sel_model)
                        else:
                            st.session_state.pop("gradcam_path", None)
                            st.session_state.pop("lime_path", None)
                        _save_prediction(result, input_image, "canvas_single_digit", xai_on, user)
                    except Exception as e:
                        st.error(f"Single digit prediction failed: {e}")

            result = st.session_state.get("last_canvas_single_result")
            if result:
                _render_results(result)
            else:
                st.markdown('<div class="result-placeholder">👈 Draw one clear digit, select model, then click <b>Predict Single Digit</b>. The result will appear here.</div>', unsafe_allow_html=True)

    # ──────────────────────────────────────────────────────────────────────
    # 2. MULTI-DIGIT CANVAS
    # ──────────────────────────────────────────────────────────────────────
    elif prediction_type == "Multi-Digit":
        st.markdown("### 🔢 Multi-Digit Prediction from Canvas")
        _clear_section("Use this for full numbers", "Example: draw `12345` or `9876`. Leave small gaps between digits so the detector can separate them.")
        col_left, col_right = st.columns([1, 1], gap="large")

        with col_left:
            brush = st.slider("Brush size", 6, 28, 14, key="multi_digit_canvas_brush")
            canvas_result = st_canvas(
                fill_color="rgba(255,255,255,0)",
                stroke_width=brush,
                stroke_color="#000000",
                background_color="#FFFFFF",
                height=220,
                width=700,
                drawing_mode="freedraw",
                key="multi_digit_canvas",
            )

            input_image = None
            if canvas_result.image_data is not None:
                arr = canvas_result.image_data.astype(np.uint8)
                if _canvas_has_ink(arr, background="white"):
                    input_image = arr[:, :, :3]

            default_index = 0
            if "cnn_deep" in AVAILABLE_MODELS:
                default_index = AVAILABLE_MODELS.index("cnn_deep")
            elif "cnn_medium" in AVAILABLE_MODELS:
                default_index = AVAILABLE_MODELS.index("cnn_medium")

            selected_model = st.selectbox(
                "Select Model",
                AVAILABLE_MODELS,
                format_func=lambda m: MODEL_LABELS.get(m, m),
                index=default_index,
                key="multi_digit_canvas_model",
            )
            confidence_threshold = st.slider("Confidence Threshold", 0, 100, 20, 5, key="multi_digit_canvas_confidence")
            max_digits = st.slider("Maximum Digits", 1, 100, 50, key="multi_digit_canvas_max_digits")
            show_debug = st.checkbox("Show mask and crops", value=True, key="multi_digit_canvas_debug")
            detect_btn = st.button(
                "🔍 Detect Multi-Digit",
                type="primary",
                use_container_width=True,
                disabled=(input_image is None),
                key="multi_digit_canvas_btn",
            )

        with col_right:
            if detect_btn and input_image is not None:
                valid, msg = _validate_canvas_digit_like(input_image, mode="multi_digit")
                if not valid:
                    st.warning(f"⚠️ {msg}")
                    return

                try:
                    from preprocessing.universal_digit_detector import detect_digits_universal
                    from backend.services.prediction_service import registry
                    model = registry.get(selected_model)
                    if model is None:
                        st.error(f"Model not found: {selected_model}")
                        return
                    with st.spinner("Detecting digits from canvas..."):
                        result = detect_digits_universal(
                            input_image,
                            model=model,
                            confidence_threshold=float(confidence_threshold),
                            max_digits=int(max_digits),
                        )
                    if not result.get("success"):
                        st.error(result.get("message", "No digits detected."))
                        if show_debug and result.get("mask_image") is not None:
                            st.image(result["mask_image"], caption="Combined Mask", use_container_width=True, clamp=True)
                        return
                    st.session_state["last_canvas_multi_digit_result"] = result
                    # ── Save to universal history ──────────────────────
                    _save_universal_prediction(
                        category="multidigit",
                        user=user,
                        input_source="canvas",
                        model_used=selected_model,
                        multidigit_result=str(result.get("number_prediction") or result.get("prediction", "")),
                        digit_count=result.get("digit_count", 0),
                        image_arr=input_image,
                        extra_data={"rows": result.get("number_rows", [])},
                    )
                    _render_universal_digit_result(result, show_debug=show_debug)
                except ModuleNotFoundError as e:
                    st.error("Missing detector file: preprocessing/universal_digit_detector.py")
                    st.code(str(e))
                except Exception as e:
                    st.error(f"Canvas multi-digit prediction failed: {e}")
            else:
                last = st.session_state.get("last_canvas_multi_digit_result")
                if last:
                    _render_universal_digit_result(last, show_debug=show_debug)
                else:
                    st.markdown('<div class="result-placeholder">👈 Draw a number sequence, keep gaps between digits, then click <b>Detect Multi-Digit</b>. The detected boxes and final number will appear here.</div>', unsafe_allow_html=True)

    # ──────────────────────────────────────────────────────────────────────
    # 3. OCR CANVAS
    # ──────────────────────────────────────────────────────────────────────
    elif prediction_type == "OCR Text / Alphabet":
        st.markdown("### 🔤 OCR Text / Alphabet Prediction from Canvas")
        _clear_section("Use this for alphabets or words", "Example: draw `A`, `HELLO`, or short handwritten words. For long paragraphs, use the upload page with document OCR.")
        from backend.services.ocr_prediction_service import (
            get_available_ocr_models,
            load_ocr_model_from_path,
            predict_full_image_ocr,
        )

        col_left, col_right = st.columns([1, 1], gap="large")
        with col_left:
            st.markdown('<div class="section-title">Draw alphabet / word / text</div>', unsafe_allow_html=True)
            ocr_brush = st.slider("Brush size", 4, 24, 10, key="ocr_canvas_brush")
            ocr_canvas = st_canvas(
                fill_color="rgba(255,255,255,0)",
                stroke_width=ocr_brush,
                stroke_color="#000000",
                background_color="#FFFFFF",
                height=160,
                width=700,
                drawing_mode="freedraw",
                key="ocr_text_canvas",
            )

            ocr_input_image = None
            if ocr_canvas.image_data is not None:
                arr = ocr_canvas.image_data.astype(np.uint8)
                if _canvas_has_ink(arr, background="white"):
                    ocr_input_image = Image.fromarray(arr).convert("RGB")

            ocr_model_map = get_available_ocr_models()
            ocr_labels = list(ocr_model_map.keys())
            if ocr_labels:
                selected_ocr_label = st.selectbox("OCR Model", ocr_labels, index=0, key="ocr_canvas_model_select")
                ocr_model_path = str(ocr_model_map[selected_ocr_label])
            else:
                ocr_model_path = None
                st.warning("No OCR model found in models/. Please add ocr_prediction_model.keras.")

            ocr_mode_label = st.selectbox(
                "OCR Reading Mode",
                ["Words or Line", "Auto Detect", "Separate Characters", "Full Image"],
                index=1,
                key="ocr_canvas_mode",
            )
            ocr_mode_map = {
                "Words or Line": "line",
                "Auto Detect": "auto",
                "Separate Characters": "character_grid",
                "Full Image": "whole",
            }
            force_upper = st.checkbox("Convert result to UPPERCASE", value=False, key="ocr_canvas_uppercase")
            ocr_btn = st.button(
                "📝 Predict OCR Text",
                type="primary",
                use_container_width=True,
                disabled=(ocr_input_image is None),
                key="ocr_canvas_btn",
            )

        with col_right:
            if ocr_btn and ocr_input_image is not None:
                with st.spinner("Running OCR from canvas..."):
                    try:
                        ocr_mode = ocr_mode_map.get(ocr_mode_label, "line")
                        external_ocr_modes = {"printed_paragraph"}
                        loaded_ocr_model = None if ocr_mode in external_ocr_modes else (load_ocr_model_from_path(ocr_model_path) if ocr_model_path else None)
                        ocr_result = predict_full_image_ocr(ocr_input_image, model=loaded_ocr_model, mode=ocr_mode)
                        if force_upper:
                            _uppercase_ocr_result_inplace(ocr_result)
                        st.session_state["last_canvas_ocr_result"] = ocr_result
                        # ── Save to universal history ──────────────────────
                        _raw_ocr_conf = ocr_result.get("confidence", 0.0) or 0.0
                        try:
                            _raw_ocr_conf = float(_raw_ocr_conf)
                            if _raw_ocr_conf > 1.0:
                                _raw_ocr_conf = _raw_ocr_conf / 100.0
                        except Exception:
                            _raw_ocr_conf = 0.0
                        _ocr_lines = ocr_result.get("lines", [])
                        _save_universal_prediction(
                            category="alphabet",
                            user=user,
                            input_source="canvas",
                            model_used=selected_ocr_label if ocr_labels else "external",
                            ocr_mode=ocr_mode,
                            ocr_text=str(ocr_result.get("text", "") or ""),
                            ocr_confidence=_raw_ocr_conf,
                            line_count=len(_ocr_lines),
                            image_arr=np.array(ocr_input_image),
                            extra_data={"segmentation_mode": ocr_result.get("segmentation_mode")},
                        )
                    except Exception as e:
                        st.error(f"Canvas OCR failed: {e}")

            ocr_result = st.session_state.get("last_canvas_ocr_result")
            if ocr_result:
                _render_ocr_result_block(ocr_result)
            else:
                st.markdown('<div class="result-placeholder">👈 Draw an alphabet, word, or short text, then click <b>Predict OCR Text</b>. OCR output and confidence will appear here.</div>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: UPLOAD PREDICTION - MERGED UPLOAD FEATURES
# ══════════════════════════════════════════════════════════════════════════════
def page_upload_prediction(user):
    st.markdown("""
    <div class="page-hero">
      <span class="mode-chip">📁 Upload</span>
      <h2>Upload Prediction</h2>
      <p>Upload an image, batch, PDF, Word, or PowerPoint file.</p>
    </div>
    """, unsafe_allow_html=True)

    upload_type = _compact_choice(
        "Choose upload type",
        ["Single Digit Image", "Multi-Digit Image", "Batch Digit Images", "OCR / PDF / Word / PowerPoint"],
        default="Single Digit Image",
        key="upload_prediction_type",
    )

    st.divider()

    # ──────────────────────────────────────────────────────────────────────
    # 1. SINGLE DIGIT IMAGE UPLOAD
    # ──────────────────────────────────────────────────────────────────────
    if upload_type == "Single Digit Image":
        st.markdown("### 🔢 Single Digit Prediction from Image")
        _clear_section("Upload one digit image", "Use this when your image contains exactly one digit. For numbers like 12345, choose Multi-Digit Image.")
        col_left, col_right = st.columns([1, 1], gap="large")

        with col_left:
            uploaded = st.file_uploader(
                "Upload single digit image",
                type=["png", "jpg", "jpeg", "bmp", "webp"],
                key="single_digit_upload_file",
            )
            input_image = None
            if uploaded is not None:
                pil_img = ImageOps.exif_transpose(Image.open(uploaded)).convert("RGB")
                st.image(pil_img, caption="Uploaded Image", use_container_width=True)
                input_image = np.array(pil_img)[:, :, ::-1]

            sel_model = st.selectbox(
                "Model",
                AVAILABLE_MODELS,
                format_func=lambda m: MODEL_LABELS.get(m, m),
                index=min(1, len(AVAILABLE_MODELS) - 1),
                key="single_digit_upload_model",
            )
            xai_on = st.checkbox("Show XAI heatmaps", value=True, key="single_digit_upload_xai")
            pred_btn = st.button(
                "🔮 Predict Single Digit",
                type="primary",
                use_container_width=True,
                disabled=(input_image is None),
                key="single_digit_upload_btn",
            )

        with col_right:
            if pred_btn and input_image is not None:
                with st.spinner("Running single digit prediction..."):
                    try:
                        pfn, gcfn, limefn, fgsmfn, reg = load_prediction_service()
                        result = pfn(input_image, sel_model)
                        st.session_state.update({
                            "last_result": result,
                            "last_upload_single_result": result,
                            "last_image": input_image,
                            "last_model": sel_model,
                            "last_input_type": "upload_single_digit",
                            "xai_on": xai_on,
                        })
                        if xai_on:
                            st.session_state["gradcam_path"] = gcfn(input_image, sel_model)
                            st.session_state["lime_path"] = limefn(input_image, sel_model)
                        else:
                            st.session_state.pop("gradcam_path", None)
                            st.session_state.pop("lime_path", None)
                        _save_prediction(result, input_image, "upload_single_digit", xai_on, user)
                    except Exception as e:
                        st.error(f"Single digit prediction failed: {e}")

            result = st.session_state.get("last_upload_single_result")
            if result:
                _render_results(result)
            else:
                st.markdown('<div class="result-placeholder">👈 Upload one clear digit image, select model, then click <b>Predict Single Digit</b>. The result will appear here.</div>', unsafe_allow_html=True)

    # ──────────────────────────────────────────────────────────────────────
    # 2. MULTI-DIGIT IMAGE UPLOAD
    # ──────────────────────────────────────────────────────────────────────
    elif upload_type == "Multi-Digit Image":
        st.markdown("### 🔢 Multi-Digit Image Prediction")
        _clear_section("Upload one image with many digits", "Use this for full numbers, forms, multiple rows, or images containing digit sequences.")
        page_multi_digit_predict(user)

    # ──────────────────────────────────────────────────────────────────────
    # 3. BATCH DIGIT IMAGE UPLOAD
    # ──────────────────────────────────────────────────────────────────────
    elif upload_type == "Batch Digit Images":
        st.markdown("### 📦 Batch Digit Image Prediction")
        _clear_section("Upload many single-digit images", "Use this when every file contains one digit and you want all predictions together in one table/report.")
        page_batch_predict(user)

    # ──────────────────────────────────────────────────────────────────────
    # 4. OCR / PDF / WORD / PPT UPLOAD
    # ──────────────────────────────────────────────────────────────────────
    elif upload_type == "OCR / PDF / Word / PowerPoint":
        st.markdown("### 🔤 OCR / Document Prediction")
        _clear_section("Upload text or documents", "Use this for image OCR, PDF pages, Word DOCX text, PowerPoint PPTX text, and embedded images.")
        page_ocr_predict(user)


def page_admin_test_prediction(user):
    st.markdown("## 🧪 Test Prediction")
    st.caption(
        "Admin testing page for canvas prediction, upload prediction, multi-digit detection, OCR, PDF, Word, and PPT support."
    )

    test_mode = _compact_choice(
        "Select test mode",
        ["🖊️ Canvas Prediction", "📁 Upload Prediction"],
        default="🖊️ Canvas Prediction",
        key="admin_test_prediction_mode",
    )

    st.divider()

    if test_mode == "🖊️ Canvas Prediction":
        page_canvas_prediction(user)
    else:
        page_upload_prediction(user)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN  –  navigation router
# ══════════════════════════════════════════════════════════════════════════════
def main():
    user = st.session_state.get("auth_user")

    # ── Not logged in → public pages only ────────────────────────────────────
    if not user:
        nav = st.session_state.get("nav", "landing")
        if nav == "login":
            page_login()
        elif nav == "register":
            page_register()
        elif nav == "forgot_email":
            page_forgot_email()
        elif nav == "forgot_otp":
            page_forgot_otp()
        elif nav == "forgot_reset":
            page_forgot_reset()
        else:
            page_landing()
        st.stop()
        return

    # ── Authenticated ────────────────────────────────────────────────────────
    _render_navbar(user)

    # CLIENT: prediction + own history only
    if user["role"] == "client":
        t1, t2, t3 = st.tabs([
            "🖊️ Canvas Prediction",
            "📁 Upload Prediction",
            "🕐 My History",
        ])

        with t1:
            page_canvas_prediction(user)
        with t2:
            page_upload_prediction(user)
        with t3:
            page_client_history(user)

    # ADMIN: monitoring + management + testing
    elif user["role"] == "admin":
        t1, t2, t3, t4, t5 = st.tabs([
            "📊 Admin Dashboard",
            "🕐 All Prediction History",
            "👥 User Management",
            "📈 Model Accuracy",
            "🧪 Test Prediction",
        ])

        with t1:
            page_admin_dashboard()
        with t2:
            page_admin_history()
        with t3:
            page_user_management()
        with t4:
            page_accuracy_admin()
        with t5:
            page_admin_test_prediction(user)

    else:
        st.error("Invalid user role. Please contact admin.")


if __name__ == "__main__":
    main()
